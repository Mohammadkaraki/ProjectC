"""
Translate ALL slides in a PowerPoint presentation
OPTIMIZED: Uses parallel processing for maximum speed
"""
import sys
import os
from pptx import Presentation
import shutil
from concurrent.futures import ThreadPoolExecutor, as_completed
import time

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from modules.slide_parser import extract_slide_structure
from modules.context_builder import build_context_map
from modules.llm_translator import translate_with_openai
from modules.rtl_converter import flip_to_rtl_layout, group_chart_elements
from modules.text_replacer import replace_text_in_slide
from modules.chart_translator import translate_charts_in_pptx
from modules.chart_collision_fixer import fix_chart_collisions_option_c
from modules.layout_translator import translate_slide_layouts
from config import Config
from utils.logger import setup_logger
from pptx import Presentation

logger = setup_logger(__name__)

def process_single_slide(input_path: str, slide_idx: int, slide_count: int):
    """
    Process a single slide: Parse → Context → Translate
    This function can run in parallel for multiple slides
    """
    logger.info(f"[Slide {slide_idx+1}/{slide_count}] Starting parallel processing...")

    start_time = time.time()

    # Parse slide structure
    structure = extract_slide_structure(input_path, slide_idx)
    logger.info(f"[Slide {slide_idx+1}/{slide_count}] Parsed {len(structure.get('elements', []))} elements")

    # Build context map
    context = build_context_map(structure)

    # Translate with OpenAI
    translations = translate_with_openai(structure, context,
                                        Config.SOURCE_LANGUAGE,
                                        Config.TARGET_LANGUAGE)

    elapsed = time.time() - start_time
    logger.info(f"[Slide {slide_idx+1}/{slide_count}] ✓ Completed in {elapsed:.1f}s")

    return {
        'slide_idx': slide_idx,
        'structure': structure,
        'translations': translations
    }

def translate_all_slides(input_path: str, output_path: str):
    """
    Translate all slides in presentation using PARALLEL processing
    Significantly faster for presentations with multiple slides
    """

    logger.info("="*60)
    logger.info("⚡ PARALLEL Translation Mode - ALL slides")
    logger.info("="*60)

    overall_start = time.time()

    # Get slide count
    prs = Presentation(input_path)
    slide_count = len(prs.slides)
    logger.info(f"Found {slide_count} slides")

    # Step 1-3: Parse, build context, translate ALL slides IN PARALLEL
    logger.info(f"\n🚀 Processing all {slide_count} slides in PARALLEL...")

    all_slides_data = [None] * slide_count  # Pre-allocate list

    # Use ThreadPoolExecutor for parallel processing
    max_workers = min(slide_count, 5)  # Max 5 parallel threads
    logger.info(f"Using {max_workers} parallel workers")

    with ThreadPoolExecutor(max_workers=max_workers) as executor:
        # Submit all slides for parallel processing
        future_to_slide = {
            executor.submit(process_single_slide, input_path, slide_idx, slide_count): slide_idx
            for slide_idx in range(slide_count)
        }

        # Collect results as they complete
        for future in as_completed(future_to_slide):
            slide_data = future.result()
            slide_idx = slide_data['slide_idx']
            all_slides_data[slide_idx] = slide_data

    parallel_time = time.time() - overall_start
    logger.info(f"\n✓ All {slide_count} slides processed in {parallel_time:.1f}s (parallel)")

    # Step 4: Group chart elements (before RTL conversion)
    logger.info("\nGrouping chart-related elements on chart slides...")
    prs = Presentation(input_path)
    total_groups = 0
    for slide_idx, slide in enumerate(prs.slides):
        groups_created = group_chart_elements(slide)
        if groups_created > 0:
            logger.info(f"  Slide {slide_idx + 1}: Created {groups_created} group(s)")
            total_groups += groups_created

    # Save the presentation with grouped charts
    grouped_temp = output_path.replace('.pptx', '_grouped_temp.pptx')
    prs.save(grouped_temp)
    logger.info(f"Created {total_groups} total chart group(s)")

    # Step 5: Convert to RTL (all slides - normal conversion)
    logger.info("\nConverting to RTL layout...")
    rtl_temp = output_path.replace('.pptx', '_rtl_temp.pptx')
    flip_to_rtl_layout(grouped_temp, rtl_temp)

    # Step 5.5: Fix chart collisions (shift charts to avoid objects)
    logger.info("\nFixing chart-to-object collisions...")
    prs_rtl = Presentation(rtl_temp)
    slide_width = prs_rtl.slide_width
    fix_chart_collisions_option_c(prs_rtl, slide_width)
    prs_rtl.save(rtl_temp)
    logger.info("Chart collision fixes applied")

    # Step 6: Replace text in ALL slides
    logger.info("\nReplacing text in all slides...")
    current_file = rtl_temp

    for slide_idx, slide_data in enumerate(all_slides_data):
        logger.info(f"[Slide {slide_idx+1}/{slide_count}] Replacing text...")

        # For last slide, output to final path
        out_file = output_path if slide_idx == slide_count-1 else current_file

        replace_text_in_slide(
            current_file,
            slide_data['translations'],
            slide_data['structure'],
            out_file,
            slide_idx
        )

    # Step 7: Translate charts
    logger.info("\nTranslating chart text...")
    chart_temp = output_path.replace('.pptx', '_chart_translated.pptx')
    translate_charts_in_pptx(output_path, chart_temp,
                            Config.SOURCE_LANGUAGE,
                            Config.TARGET_LANGUAGE)
    if os.path.exists(chart_temp):
        os.remove(output_path)
        os.rename(chart_temp, output_path)

    # Step 8: Translate layouts
    logger.info("\nTranslating layout backgrounds...")
    layout_out = output_path.replace('.pptx', '_with_layouts.pptx')
    translate_slide_layouts(output_path, layout_out)

    if os.path.exists(layout_out):
        os.remove(output_path)
        os.rename(layout_out, output_path)

    # Cleanup
    if os.path.exists(grouped_temp):
        os.remove(grouped_temp)
    if os.path.exists(rtl_temp):
        os.remove(rtl_temp)

    total_time = time.time() - overall_start
    logger.info("="*60)
    logger.info(f"✓ SUCCESS! Translated {slide_count} slides")
    logger.info(f"⚡ Total time: {total_time:.1f}s (with parallel processing)")
    logger.info(f"📊 Parallel phase: {parallel_time:.1f}s for {slide_count} slides")
    logger.info(f"💾 Output: {output_path}")
    logger.info("="*60)

    return output_path

if __name__ == "__main__":
    if len(sys.argv) > 2:
        translate_all_slides(sys.argv[1], sys.argv[2])
    else:
        print("Usage: python translate_all_slides.py <input.pptx> <output.pptx>")
