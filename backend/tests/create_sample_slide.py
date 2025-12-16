"""
Create a sample consulting slide for testing
This script generates a PowerPoint slide with:
- Title
- Header/key message
- Multiple bullet points
- Sub-bullets
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import sys
import os

# Add parent directory to path
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

def create_sample_slide(output_path: str):
    """
    Create a sample consulting slide

    Args:
        output_path: Path to save the sample .pptx file
    """
    # Create presentation
    prs = Presentation()

    # Set slide dimensions (16:9 aspect ratio)
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Add a blank slide
    blank_slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_slide_layout)

    # Add Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5),  # left, top
        Inches(9), Inches(0.8)      # width, height
    )
    title_frame = title_box.text_frame
    title_frame.text = "Market Expansion Strategy"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.name = "Calibri"

    # Add Header/Key Message
    header_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.5),  # left, top
        Inches(9), Inches(0.6)      # width, height
    )
    header_frame = header_box.text_frame
    header_frame.text = "Key findings from MENA region analysis"
    header_para = header_frame.paragraphs[0]
    header_para.font.size = Pt(24)
    header_para.font.name = "Calibri"
    header_para.font.color.rgb = RGBColor(0, 51, 102)  # Dark blue

    # Add Bullet Points
    bullet_box = slide.shapes.add_textbox(
        Inches(1), Inches(2.5),   # left, top
        Inches(8), Inches(4)       # width, height
    )
    bullet_frame = bullet_box.text_frame
    bullet_frame.word_wrap = True

    # Bullet 1
    p1 = bullet_frame.paragraphs[0]
    p1.text = "Revenue growth projected at 15% annually through 2027"
    p1.level = 0
    p1.font.size = Pt(18)
    p1.font.name = "Calibri"

    # Bullet 2
    p2 = bullet_frame.add_paragraph()
    p2.text = "Market share expansion opportunities in UAE and Saudi Arabia"
    p2.level = 0
    p2.font.size = Pt(18)
    p2.font.name = "Calibri"

    # Sub-bullet 2.1
    p2_1 = bullet_frame.add_paragraph()
    p2_1.text = "Digital transformation initiatives driving adoption"
    p2_1.level = 1
    p2_1.font.size = Pt(16)
    p2_1.font.name = "Calibri"

    # Bullet 3
    p3 = bullet_frame.add_paragraph()
    p3.text = "Competitive landscape analysis reveals strategic positioning advantages"
    p3.level = 0
    p3.font.size = Pt(18)
    p3.font.name = "Calibri"

    # Bullet 4
    p4 = bullet_frame.add_paragraph()
    p4.text = "Recommended investment: $2.5M for market entry phase"
    p4.level = 0
    p4.font.size = Pt(18)
    p4.font.name = "Calibri"

    # Save presentation
    prs.save(output_path)
    print(f"Sample slide created: {output_path}")
    print("\nSlide contents:")
    print("- Title: Market Expansion Strategy")
    print("- Header: Key findings from MENA region analysis")
    print("- 4 main bullet points + 1 sub-bullet")

if __name__ == "__main__":
    # Create sample slide
    output_file = os.path.join(
        os.path.dirname(os.path.dirname(__file__)),
        "tests", "fixtures", "sample_slide.pptx"
    )

    # Ensure fixtures directory exists
    os.makedirs(os.path.dirname(output_file), exist_ok=True)

    create_sample_slide(output_file)
