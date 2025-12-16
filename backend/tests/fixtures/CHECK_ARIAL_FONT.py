"""Check font in Template_ARIAL_CLASSIC.pptx"""
from pptx import Presentation

prs = Presentation('Template_ARIAL_CLASSIC.pptx')

output_lines = []
output_lines.append("Checking fonts in Template_ARIAL_CLASSIC.pptx\n")

for slide_idx, slide in enumerate(prs.slides):
    output_lines.append(f"Slide {slide_idx + 1}:")
    shape_count = 0

    for shape in slide.shapes:
        if hasattr(shape, 'text_frame') and shape.text_frame:
            text = shape.text_frame.text.strip()
            if text:
                shape_count += 1
                # Get first run with text
                font_name = "None"
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.text.strip():
                            font_name = run.font.name if run.font.name else "None"
                            break
                    if font_name != "None":
                        break

                text_preview = text[:30] + '...' if len(text) > 30 else text
                output_lines.append(f"  Shape {shape_count}: Font={font_name}, Text='{text_preview}'")

    output_lines.append(f"  Total text shapes: {shape_count}\n")

# Write to file
with open('ARIAL_FONT_CHECK.txt', 'w', encoding='utf-8') as f:
    f.write('\n'.join(output_lines))

print("Font check complete. Results saved to ARIAL_FONT_CHECK.txt")
