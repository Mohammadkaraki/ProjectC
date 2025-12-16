"""
Check the translated slide for Arabic text and formatting
"""
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
import json

def check_slide(pptx_path):
    """Check the translated slide"""
    prs = Presentation(pptx_path)
    slide = prs.slides[0]

    results = {
        "total_shapes": len(slide.shapes),
        "text_shapes": [],
        "arabic_count": 0,
        "right_aligned_count": 0,
        "formatting_preserved": 0
    }

    for idx, shape in enumerate(slide.shapes):
        if shape.has_text_frame:
            text = shape.text.strip()
            if text and shape.text_frame.paragraphs:
                # Check if it's Arabic
                has_arabic = any('\u0600' <= c <= '\u06FF' for c in text)

                # Get formatting
                para = shape.text_frame.paragraphs[0]
                is_right_aligned = para.alignment == PP_ALIGN.RIGHT

                shape_info = {
                    "index": idx,
                    "text_preview": text[:40] if not has_arabic else "[Arabic text]",
                    "has_arabic": has_arabic,
                    "right_aligned": is_right_aligned,
                    "font": None,
                    "size": None,
                    "bold": None
                }

                if para.runs:
                    first_run = para.runs[0]
                    shape_info["font"] = first_run.font.name
                    shape_info["size"] = str(first_run.font.size) if first_run.font.size else "None"
                    shape_info["bold"] = first_run.font.bold

                results["text_shapes"].append(shape_info)

                if has_arabic:
                    results["arabic_count"] += 1
                if is_right_aligned:
                    results["right_aligned_count"] += 1
                if shape_info["font"] and shape_info["size"] != "None":
                    results["formatting_preserved"] += 1

    # Print results
    print(json.dumps(results, indent=2, ensure_ascii=False))

    # Summary
    print(f"\n{'='*60}")
    print(f"SUMMARY:")
    print(f"  Total shapes: {results['total_shapes']}")
    print(f"  Text shapes: {len(results['text_shapes'])}")
    print(f"  Arabic shapes: {results['arabic_count']}")
    print(f"  Right-aligned: {results['right_aligned_count']}")
    print(f"  Formatting preserved: {results['formatting_preserved']}")
    print(f"{'='*60}")

if __name__ == "__main__":
    import sys
    if len(sys.argv) > 1:
        check_slide(sys.argv[1])
    else:
        print("Usage: python check_result.py <pptx_file>")
