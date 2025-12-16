"""
Create a LEFT-ALIGNED sample slide to test RTL flipping
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def create_left_aligned_slide(output_path: str):
    """Create a slide with LEFT-aligned content"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    # Title - LEFT side only
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5),  # LEFT side
        Inches(4), Inches(0.8)      # Only 4 inches wide (left half)
    )
    title_frame = title_box.text_frame
    title_frame.text = "LEFT-ALIGNED Title"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True

    # Header - LEFT side
    header_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.5),
        Inches(4.5), Inches(0.6)
    )
    header_frame = header_box.text_frame
    header_frame.text = "This text is on the LEFT side"
    header_para = header_frame.paragraphs[0]
    header_para.font.size = Pt(20)

    # Bullet box - LEFT side
    bullet_box = slide.shapes.add_textbox(
        Inches(1), Inches(2.5),
        Inches(4), Inches(3)
    )
    bullet_frame = bullet_box.text_frame

    p1 = bullet_frame.paragraphs[0]
    p1.text = "First bullet on LEFT"
    p1.font.size = Pt(16)

    p2 = bullet_frame.add_paragraph()
    p2.text = "Second bullet on LEFT"
    p2.font.size = Pt(16)

    p3 = bullet_frame.add_paragraph()
    p3.text = "Third bullet on LEFT"
    p3.font.size = Pt(16)

    prs.save(output_path)
    print(f"Left-aligned slide created: {output_path}")

if __name__ == "__main__":
    import os
    output_file = os.path.join(
        os.path.dirname(os.path.dirname(__file__)),
        "tests", "fixtures", "left_aligned_slide.pptx"
    )
    create_left_aligned_slide(output_file)
