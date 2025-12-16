"""
Check XML for text direction and bold formatting
"""
from pptx import Presentation
import zipfile
from lxml import etree

def check_file(pptx_path):
    """Check both PowerPoint API and raw XML"""
    print(f"\n{'='*80}")
    print(f"CHECKING: {pptx_path}")
    print(f"{'='*80}\n")

    # Check via PowerPoint API
    print("=== VIA POWERPOINT API ===")
    prs = Presentation(pptx_path)
    slide = prs.slides[0]

    for idx, shape in enumerate(slide.shapes):
        if shape.has_text_frame and shape.text.strip():
            text = shape.text.strip()
            has_arabic = any('\u0600' <= c <= '\u06FF' for c in text)

            if has_arabic:
                print(f"\nShape {idx} (Arabic):")
                print(f"  Text (first 50 chars): {text[:50]}")
                print(f"  Text (Unicode escape): {repr(text[:30])}")

                # Check formatting
                for para_idx, para in enumerate(shape.text_frame.paragraphs):
                    if para.runs:
                        run = para.runs[0]
                        print(f"  Para {para_idx}: Font={run.font.name}, Size={run.font.size}, Bold={run.font.bold}")
                break

    # Check raw XML
    print("\n=== VIA RAW XML ===")
    with zipfile.ZipFile(pptx_path, 'r') as zip_ref:
        # Read slide1.xml
        slide_xml = zip_ref.read('ppt/slides/slide1.xml')
        root = etree.fromstring(slide_xml)

        # Namespaces
        ns = {
            'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
            'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'
        }

        # Find all text runs
        text_runs = root.xpath('.//a:t', namespaces=ns)

        for idx, t_elem in enumerate(text_runs[:3]):  # First 3
            text = t_elem.text
            if text:
                has_arabic = any('\u0600' <= c <= '\u06FF' for c in text)
                if has_arabic:
                    print(f"\nText run {idx}:")
                    print(f"  Text: {text[:50]}")
                    print(f"  Unicode: {repr(text[:30])}")

                    # Check parent paragraph properties for RTL
                    para = t_elem.getparent().getparent()  # r -> p
                    pPr = para.find('.//a:pPr', namespaces=ns)
                    if pPr is not None:
                        rtl_attr = pPr.get('rtl')
                        print(f"  RTL attribute: {rtl_attr}")

                    # Check run properties for bold
                    rPr = t_elem.getparent().find('.//a:rPr', namespaces=ns)
                    if rPr is not None:
                        bold_attr = rPr.get('b')
                        print(f"  Bold attribute: {bold_attr}")

                    break

if __name__ == "__main__":
    import sys
    if len(sys.argv) > 1:
        check_file(sys.argv[1])
    else:
        print("Usage: python check_xml_and_bold.py <pptx_file>")
