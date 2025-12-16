"""
Quick verification script to check if the fixes worked
"""
from pptx import Presentation
import sys

def verify_translation(pptx_path):
    """Verify the translated slide"""
    print(f"\n{'='*80}")
    print(f"VERIFYING: {pptx_path}")
    print(f"{'='*80}\n")

    prs = Presentation(pptx_path)
    slide = prs.slides[0]

    print(f"Total shapes: {len(slide.shapes)}\n")

    shape_count = 0
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text.strip()
            if text:
                shape_count += 1

                # Check if it's Arabic text (contains Arabic characters)
                has_arabic = any('\u0600' <= c <= '\u06FF' for c in text)

                # Get first run to check formatting
                if shape.text_frame.paragraphs and shape.text_frame.paragraphs[0].runs:
                    first_run = shape.text_frame.paragraphs[0].runs[0]
                    font_name = first_run.font.name
                    font_size = first_run.font.size
                    is_bold = first_run.font.bold

                    # Check paragraph alignment
                    from pptx.enum.text import PP_ALIGN
                    para = shape.text_frame.paragraphs[0]
                    is_right_aligned = para.alignment == PP_ALIGN.RIGHT

                    print(f"Shape {shape_count}:")
                    print(f"  Text: {text[:60]}{'...' if len(text) > 60 else ''}")
                    print(f"  Arabic: {has_arabic}")
                    print(f"  Font: {font_name}")
                    print(f"  Size: {font_size}")
                    print(f"  Bold: {is_bold}")
                    print(f"  Right-aligned: {is_right_aligned}")
                    print()

    print(f"{'='*80}")
    print(f"Total text shapes processed: {shape_count}")
    print(f"{'='*80}\n")

if __name__ == "__main__":
    if len(sys.argv) > 1:
        verify_translation(sys.argv[1])
    else:
        print("Usage: python verify_fixes.py <pptx_file>")
