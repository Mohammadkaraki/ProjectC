"""
Compare original vs translated slide formatting
"""
from pptx import Presentation

def compare_slides(original_path, translated_path):
    """Compare original and translated slides"""
    print("\n" + "="*80)
    print("FORMATTING COMPARISON")
    print("="*80 + "\n")

    orig_prs = Presentation(original_path)
    trans_prs = Presentation(translated_path)

    orig_slide = orig_prs.slides[0]
    trans_slide = trans_prs.slides[0]

    print(f"Original shapes: {len(orig_slide.shapes)}")
    print(f"Translated shapes: {len(trans_slide.shapes)}\n")

    # Compare text shapes
    orig_shapes = []
    for shape in orig_slide.shapes:
        if shape.has_text_frame and shape.text.strip():
            if shape.text_frame.paragraphs and shape.text_frame.paragraphs[0].runs:
                para = shape.text_frame.paragraphs[0]
                run = para.runs[0]
                orig_shapes.append({
                    'text': shape.text.strip()[:30],
                    'font': run.font.name,
                    'size': run.font.size,
                    'bold': run.font.bold
                })

    trans_shapes = []
    for shape in trans_slide.shapes:
        if shape.has_text_frame and shape.text.strip():
            text = shape.text.strip()
            has_arabic = any('\u0600' <= c <= '\u06FF' for c in text)
            if shape.text_frame.paragraphs:
                para = shape.text_frame.paragraphs[0]
                if para.runs:
                    run = para.runs[0]
                    trans_shapes.append({
                        'text': '[Arabic]' if has_arabic else text[:30],
                        'font': run.font.name,
                        'size': run.font.size,
                        'bold': run.font.bold,
                        'has_runs': len(para.runs)
                    })
                else:
                    trans_shapes.append({
                        'text': '[Arabic]' if has_arabic else text[:30],
                        'font': None,
                        'size': None,
                        'bold': None,
                        'has_runs': 0
                    })

    print(f"{'ORIGINAL':<40} | {'TRANSLATED':<40}")
    print("-" * 80)
    for i in range(min(len(orig_shapes), len(trans_shapes))):
        o = orig_shapes[i]
        t = trans_shapes[i]
        o_font = o['font'] or 'None'
        o_size = str(o['size']) if o['size'] else 'None'
        o_bold = str(o['bold']) if o['bold'] is not None else 'None'
        t_font = t['font'] or 'None'
        t_size = str(t['size']) if t['size'] else 'None'
        t_bold = str(t['bold']) if t['bold'] is not None else 'None'
        t_runs = str(t.get('has_runs', '?'))
        print(f"{o_font} / {o_size} / Bold:{o_bold:<5} | {t_font} / {t_size} / Bold:{t_bold:<5} / Runs:{t_runs}")

    print("\n" + "="*80)

if __name__ == "__main__":
    import sys
    if len(sys.argv) > 2:
        compare_slides(sys.argv[1], sys.argv[2])
    else:
        print("Usage: python compare_formatting.py <original.pptx> <translated.pptx>")
