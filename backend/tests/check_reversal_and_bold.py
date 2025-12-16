"""
Check if Arabic text is reversed and if bold formatting is preserved
"""
from pptx import Presentation
import zipfile
from lxml import etree
import json

def check_text_reversal(pptx_path):
    """
    Check if Arabic text is reversed by looking at character order
    """
    results = {
        "file": pptx_path,
        "arabic_texts": [],
        "bold_status": []
    }

    prs = Presentation(pptx_path)
    slide = prs.slides[0]

    for idx, shape in enumerate(slide.shapes):
        if shape.has_text_frame and shape.text.strip():
            text = shape.text.strip()
            has_arabic = any('\u0600' <= c <= '\u06FF' for c in text)

            if has_arabic:
                # Get first and last characters to check reversal
                first_char = ord(text[0])
                last_char = ord(text[-1])

                # Check if it's in Arabic range
                is_first_arabic = 0x0600 <= first_char <= 0x06FF
                is_last_arabic = 0x0600 <= last_char <= 0x06FF

                # Get hex representation for verification
                first_10_hex = ' '.join([hex(ord(c)) for c in text[:10] if ord(c) >= 0x0600])
                last_10_hex = ' '.join([hex(ord(c)) for c in text[-10:] if ord(c) >= 0x0600])

                # Check bold
                bold_status = None
                for para in shape.text_frame.paragraphs:
                    if para.runs:
                        bold_status = para.runs[0].font.bold
                        break

                results["arabic_texts"].append({
                    "shape_idx": idx,
                    "text_length": len(text),
                    "first_char_hex": hex(first_char),
                    "last_char_hex": hex(last_char),
                    "is_first_arabic": is_first_arabic,
                    "is_last_arabic": is_last_arabic,
                    "first_10_chars": first_10_hex,
                    "last_10_chars": last_10_hex,
                    "bold": bold_status
                })

    return results

def check_original_bold(pptx_path):
    """Check bold formatting in original file"""
    results = {
        "shapes_with_bold": []
    }

    prs = Presentation(pptx_path)
    slide = prs.slides[0]

    for idx, shape in enumerate(slide.shapes):
        if shape.has_text_frame and shape.text.strip():
            for para in shape.text_frame.paragraphs:
                if para.runs:
                    if para.runs[0].font.bold:
                        results["shapes_with_bold"].append({
                            "shape_idx": idx,
                            "text_preview": shape.text.strip()[:30],
                            "bold": True
                        })
                    break

    return results

if __name__ == "__main__":
    import sys

    print("\n" + "="*80)
    print("CHECKING ORIGINAL FILE FOR BOLD")
    print("="*80)
    orig_results = check_original_bold("tests/fixtures/Template.pptx")
    print(json.dumps(orig_results, indent=2))

    print("\n" + "="*80)
    print("CHECKING TRANSLATED FILE")
    print("="*80)
    trans_results = check_text_reversal("tests/fixtures/Template_FINAL_FIXED.pptx")
    print(json.dumps(trans_results, indent=2))

    print("\n" + "="*80)
    print("ANALYSIS")
    print("="*80)

    if orig_results["shapes_with_bold"]:
        print(f"\nOriginal file has {len(orig_results['shapes_with_bold'])} shapes with BOLD text")
    else:
        print("\nOriginal file has NO bold text")

    for text_info in trans_results["arabic_texts"]:
        print(f"\nArabic text in shape {text_info['shape_idx']}:")
        print(f"  Bold status: {text_info['bold']}")
        print(f"  First char: {text_info['first_char_hex']} (Arabic: {text_info['is_first_arabic']})")
        print(f"  Last char: {text_info['last_char_hex']} (Arabic: {text_info['is_last_arabic']})")
