"""
Text Replacer Module
Replaces original English text with translated Arabic text in PowerPoint slides
"""
from pptx import Presentation
from typing import Dict, List, Any
import sys
import os

# Add parent directory to path
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from utils.logger import setup_logger

logger = setup_logger(__name__)

def replace_text_in_slide(
    pptx_path: str,
    translations: Dict[str, Any],
    slide_structure: Dict[str, Any],
    output_path: str,
    slide_index: int = 0
) -> None:
    """
    Replace original text with translated text in PowerPoint slide

    Args:
        pptx_path: Path to PowerPoint file (already RTL-converted)
        translations: Dictionary of translations from llm_translator
        slide_structure: Slide structure from slide_parser
        output_path: Path to save final output
        slide_index: Index of slide to process (default: 0)
    """
    logger.info(f"Replacing text with translations: {pptx_path}")

    try:
        # Load presentation
        prs = Presentation(pptx_path)

        # Get target slide
        if slide_index >= len(prs.slides):
            raise ValueError(f"Slide index {slide_index} out of range")

        slide = prs.slides[slide_index]

        # Create a mapping of element_id to shape object for quick lookup
        # Must handle grouped shapes recursively (same as parser logic)
        shape_map = {}
        element_counter = [0]

        def map_shapes_recursive(shape, parent_id=""):
            """Recursively map shapes including groups (MUST match parser logic)"""
            from pptx.enum.shapes import MSO_SHAPE_TYPE

            # Handle GROUPED SHAPES - recurse into nested shapes
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for sub_idx, sub_shape in enumerate(shape.shapes):
                    sub_id = f"{parent_id}_group{sub_idx}" if parent_id else f"group{sub_idx}"
                    map_shapes_recursive(sub_shape, sub_id)
                return

            # Handle TABLES - one entry per cell with text (matches parser logic)
            if hasattr(shape, 'has_table') and shape.has_table:
                has_text_cells = False
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            shape_id = f"table_{element_counter[0]}"
                            element_counter[0] += 1
                            shape_map[shape_id] = shape  # Map same shape multiple times (one per cell)
                            has_text_cells = True
                if has_text_cells:
                    return  # Table handled
                return  # No text in table, skip

            # Handle TEXT FRAMES - only count if has text
            if not shape.has_text_frame:
                return  # No text frame, skip

            text = shape.text.strip()
            if not text:
                return  # Empty text, skip

            # Shape has text - add to map
            shape_id = f"shape_{element_counter[0]}"
            element_counter[0] += 1
            shape_map[shape_id] = shape

        # Build shape map from all top-level shapes
        for shape in slide.shapes:
            map_shapes_recursive(shape)

        # Replace text for each element
        for element in slide_structure["elements"]:
            element_id = element["element_id"]

            if element_id not in translations:
                logger.warning(f"No translation found for element: {element_id}")
                continue

            if element_id not in shape_map:
                logger.warning(f"Shape not found for element: {element_id}")
                continue

            shape = shape_map[element_id]

            # Get translation
            translation = translations[element_id]

            # Replace based on element type
            if isinstance(translation, list):
                # Bullet group - replace each bullet
                _replace_bullets(shape, translation)
            else:
                # Single text element
                _replace_single_text(shape, translation)

        # Save final presentation
        prs.save(output_path)
        logger.info(f"Text replacement complete. Saved to: {output_path}")

    except Exception as e:
        logger.error(f"Error during text replacement: {str(e)}", exc_info=True)
        raise

def _replace_single_text(shape, translation: str) -> None:
    """
    Replace text in a single text element (title, header, text_box)
    PRESERVES original formatting (bold, italic, font size, color)

    Args:
        shape: PowerPoint shape object
        translation: Translated text
    """
    if not shape.has_text_frame:
        logger.warning("Shape has no text frame")
        return

    try:
        text_frame = shape.text_frame

        # Save original formatting from first paragraph/run
        original_formatting = {}
        original_para_formatting = {}

        if text_frame.paragraphs:
            first_para = text_frame.paragraphs[0]

            # Save paragraph-level formatting
            original_para_formatting = {
                'line_spacing': first_para.line_spacing,
                'space_before': first_para.space_before,
                'space_after': first_para.space_after,
            }

            # Save run-level formatting
            if first_para.runs:
                first_run = first_para.runs[0]
                original_formatting = {
                    'bold': first_run.font.bold,
                    'italic': first_run.font.italic,
                    'underline': first_run.font.underline,
                    'font_name': first_run.font.name,
                    'font_size': first_run.font.size,
                    'color': first_run.font.color.rgb if first_run.font.color and hasattr(first_run.font.color, 'rgb') else None
                }

        # Clear and add new paragraph with run
        text_frame.clear()

        # FIX: Enable word wrap to prevent text overflow
        text_frame.word_wrap = True

        # FIX: Auto-shrink text to fit shape (better for Arabic)
        from pptx.enum.text import MSO_AUTO_SIZE
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

        # FIX: Optimal margins for Arabic text (smaller = more space)
        from pptx.util import Inches
        text_frame.margin_left = Inches(0.03)
        text_frame.margin_right = Inches(0.03)
        text_frame.margin_top = Inches(0.03)
        text_frame.margin_bottom = Inches(0.03)

        paragraph = text_frame.add_paragraph()

        # Create run and set text (this ensures we have a run to format)
        run = paragraph.add_run()
        run.text = translation

        # Maintain RIGHT alignment (no RTL attribute to avoid text reversal)
        from pptx.enum.text import PP_ALIGN
        paragraph.alignment = PP_ALIGN.RIGHT

        # NOTE: We do NOT set rtl="1" because Arabic text from LLM is already
        # in correct logical order. Setting rtl="1" would reverse it again.

        # Restore paragraph-level formatting
        if original_para_formatting:
            if original_para_formatting.get('line_spacing') is not None:
                paragraph.line_spacing = original_para_formatting['line_spacing']
            if original_para_formatting.get('space_before') is not None:
                paragraph.space_before = original_para_formatting['space_before']
            if original_para_formatting.get('space_after') is not None:
                paragraph.space_after = original_para_formatting['space_after']

        # Restore run-level formatting
        if original_formatting:
            run = paragraph.runs[0]
            if original_formatting.get('bold') is not None:
                run.font.bold = original_formatting['bold']
            if original_formatting.get('italic') is not None:
                run.font.italic = original_formatting['italic']
            if original_formatting.get('underline') is not None:
                run.font.underline = original_formatting['underline']

            # Font name: Use Arial (classic standard Arabic font)
            if original_formatting.get('font_name'):
                font_name = original_formatting['font_name']
                # Common fonts that support Arabic
                arabic_fonts = ['Arial', 'Calibri', 'Times New Roman', 'Tahoma',
                               'Simplified Arabic', 'Traditional Arabic', 'Arabic Typesetting',
                               'Sakkal Majalla', 'Dubai', 'Segoe UI']
                # If original font doesn't support Arabic, use Arial
                if not any(af.lower() in font_name.lower() for af in arabic_fonts):
                    font_name = 'Arial'
                    logger.debug(f"Font {original_formatting['font_name']} doesn't support Arabic, using Arial")
                else:
                    # Even if it supports Arabic, prefer Arial for classic appearance
                    font_name = 'Arial'
                run.font.name = font_name

            # Font size: Reduce by 10% for Arabic (better fitting)
            if original_formatting.get('font_size'):
                original_size = original_formatting['font_size']
                run.font.size = int(original_size * 0.9)
            if original_formatting.get('color'):
                # FIX: Convert gray colors to black for better Arabic readability
                color = original_formatting['color']
                # Check if color is gray (common gray values)
                gray_colors = ['CCCCCC', 'CCCECE', 'C0C0C0', 'CACACA', 'D3D3D3', 'BEBEBE']
                if str(color).upper() in gray_colors:
                    from pptx.dml.color import RGBColor
                    run.font.color.rgb = RGBColor(0, 0, 0)  # Black
                    logger.debug(f"Converted gray color {color} to black for Arabic text")
                else:
                    run.font.color.rgb = original_formatting['color']

        logger.debug(f"Replaced single text: '{translation[:50]}...' (formatting preserved)")

    except Exception as e:
        logger.warning(f"Could not replace single text: {str(e)}")

def _replace_bullets(shape, translations: List[str]) -> None:
    """
    Replace bullet points with translated versions
    PRESERVES original formatting (bold, italic, font size, color) and hierarchy

    Args:
        shape: PowerPoint shape object
        translations: List of translated bullet texts
    """
    if not shape.has_text_frame:
        logger.warning("Shape has no text frame for bullets")
        return

    try:
        text_frame = shape.text_frame

        # Save original bullet levels AND formatting before clearing
        original_data = []
        for paragraph in text_frame.paragraphs:
            if paragraph.text.strip():
                para_data = {
                    'level': paragraph.level,
                    'formatting': {},
                    'para_formatting': {
                        'line_spacing': paragraph.line_spacing,
                        'space_before': paragraph.space_before,
                        'space_after': paragraph.space_after,
                    }
                }

                # Save run-level formatting from first run
                if paragraph.runs:
                    first_run = paragraph.runs[0]
                    para_data['formatting'] = {
                        'bold': first_run.font.bold,
                        'italic': first_run.font.italic,
                        'underline': first_run.font.underline,
                        'font_name': first_run.font.name,
                        'font_size': first_run.font.size,
                        'color': first_run.font.color.rgb if first_run.font.color and hasattr(first_run.font.color, 'rgb') else None
                    }

                original_data.append(para_data)

        # Ensure we have enough data (pad with defaults if needed)
        while len(original_data) < len(translations):
            original_data.append({'level': 0, 'formatting': {}, 'para_formatting': {}})

        # Clear existing content
        text_frame.clear()

        # FIX: Enable word wrap to prevent text overflow
        text_frame.word_wrap = True

        # FIX: Auto-shrink text to fit shape (better for Arabic)
        from pptx.enum.text import MSO_AUTO_SIZE
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

        # FIX: Optimal margins for Arabic text (smaller = more space)
        from pptx.util import Inches
        text_frame.margin_left = Inches(0.03)
        text_frame.margin_right = Inches(0.03)
        text_frame.margin_top = Inches(0.03)
        text_frame.margin_bottom = Inches(0.03)

        # Add translated bullets with original hierarchy and formatting
        for i, translation_text in enumerate(translations):
            paragraph = text_frame.add_paragraph()

            # Create run and set text (ensures we have a run to format)
            run = paragraph.add_run()
            run.text = translation_text

            paragraph.level = original_data[i]['level'] if i < len(original_data) else 0

            # Maintain RIGHT alignment (no RTL attribute to avoid text reversal)
            from pptx.enum.text import PP_ALIGN
            paragraph.alignment = PP_ALIGN.RIGHT

            # NOTE: We do NOT set rtl="1" because Arabic text from LLM is already
            # in correct logical order. Setting rtl="1" would reverse it again.

            # Restore paragraph-level formatting
            if original_data[i]['para_formatting']:
                para_fmt = original_data[i]['para_formatting']
                if para_fmt.get('line_spacing') is not None:
                    paragraph.line_spacing = para_fmt['line_spacing']
                if para_fmt.get('space_before') is not None:
                    paragraph.space_before = para_fmt['space_before']
                if para_fmt.get('space_after') is not None:
                    paragraph.space_after = para_fmt['space_after']

            # Restore run-level formatting
            if original_data[i]['formatting']:
                run = paragraph.runs[0]
                fmt = original_data[i]['formatting']

                if fmt.get('bold') is not None:
                    run.font.bold = fmt['bold']
                if fmt.get('italic') is not None:
                    run.font.italic = fmt['italic']
                if fmt.get('underline') is not None:
                    run.font.underline = fmt['underline']

                # Font name: Use Arial (classic standard Arabic font)
                if fmt.get('font_name'):
                    font_name = fmt['font_name']
                    # Common fonts that support Arabic
                    arabic_fonts = ['Arial', 'Calibri', 'Times New Roman', 'Tahoma',
                                   'Simplified Arabic', 'Traditional Arabic', 'Arabic Typesetting',
                                   'Sakkal Majalla', 'Dubai', 'Segoe UI']
                    # If original font doesn't support Arabic, use Arial
                    if not any(af.lower() in font_name.lower() for af in arabic_fonts):
                        font_name = 'Arial'
                        logger.debug(f"Font {fmt['font_name']} doesn't support Arabic, using Arial")
                    else:
                        # Even if it supports Arabic, prefer Arial for classic appearance
                        font_name = 'Arial'
                    run.font.name = font_name

                # Font size: Reduce by 10% for Arabic (better fitting)
                if fmt.get('font_size'):
                    original_size = fmt['font_size']
                    run.font.size = int(original_size * 0.9)
                if fmt.get('color'):
                    # FIX: Convert gray colors to black for better Arabic readability
                    color = fmt['color']
                    # Check if color is gray (common gray values)
                    gray_colors = ['CCCCCC', 'CCCECE', 'C0C0C0', 'CACACA', 'D3D3D3', 'BEBEBE']
                    if str(color).upper() in gray_colors:
                        from pptx.dml.color import RGBColor
                        run.font.color.rgb = RGBColor(0, 0, 0)  # Black
                        logger.debug(f"Converted gray color {color} to black for Arabic text")
                    else:
                        run.font.color.rgb = fmt['color']

        logger.debug(f"Replaced {len(translations)} bullets (formatting preserved)")

    except Exception as e:
        logger.warning(f"Could not replace bullets: {str(e)}")

def _set_paragraph_rtl(paragraph) -> None:
    """
    Set RTL property on paragraph via XML manipulation

    This sets the text direction to right-to-left for proper Arabic rendering.
    The rtl="1" attribute tells PowerPoint this is RTL text.

    Args:
        paragraph: PowerPoint paragraph object
    """
    try:
        from lxml import etree

        # XML namespaces for PowerPoint
        NAMESPACES = {
            'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'
        }

        # Get paragraph XML element
        p_element = paragraph._element

        # Get or create <a:pPr> (paragraph properties)
        pPr = p_element.find('.//a:pPr', namespaces=NAMESPACES)

        if pPr is None:
            # Create <a:pPr> if it doesn't exist
            pPr = etree.Element(f"{{{NAMESPACES['a']}}}pPr")
            # Insert as first child
            p_element.insert(0, pPr)

        # Set RTL attribute
        pPr.set('rtl', '1')  # 1 = RTL, 0 = LTR

        logger.debug("Set paragraph RTL property")

    except Exception as e:
        logger.warning(f"Could not set paragraph RTL: {str(e)}")

def _preserve_formatting(original_run, target_run) -> None:
    """
    Preserve font formatting from original to translated text
    (OPTIONAL - not used in MVP for simplicity)

    Args:
        original_run: Original text run
        target_run: Translated text run
    """
    try:
        if original_run.font.bold is not None:
            target_run.font.bold = original_run.font.bold

        if original_run.font.italic is not None:
            target_run.font.italic = original_run.font.italic

        if original_run.font.size is not None:
            target_run.font.size = original_run.font.size

        if original_run.font.name:
            target_run.font.name = original_run.font.name

        # Color preservation
        if original_run.font.color and original_run.font.color.rgb:
            target_run.font.color.rgb = original_run.font.color.rgb

    except Exception as e:
        logger.warning(f"Could not preserve formatting: {str(e)}")

if __name__ == "__main__":
    # Test text replacement
    if len(sys.argv) > 3:
        input_file = sys.argv[1]
        translations_json = sys.argv[2]  # JSON file with translations
        output_file = sys.argv[3]

        import json
        with open(translations_json, 'r', encoding='utf-8') as f:
            translations = json.load(f)

        from modules.slide_parser import extract_slide_structure
        structure = extract_slide_structure(input_file)

        replace_text_in_slide(input_file, translations, structure, output_file)
        print(f"Text replacement complete: {output_file}")
    else:
        print("Usage: python text_replacer.py <input.pptx> <translations.json> <output.pptx>")
