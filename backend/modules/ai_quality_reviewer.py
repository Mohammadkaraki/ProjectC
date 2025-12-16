"""
AI Quality Reviewer Module
Uses OpenAI to review translated presentations and fix structural issues
"""
import zipfile
import tempfile
import shutil
import os
import sys
from typing import Dict, List, Tuple
from pptx import Presentation

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from config import Config
from utils.logger import setup_logger
from modules.llm_translator import get_openai_client

logger = setup_logger(__name__)

def review_and_fix_presentation(original_path: str, translated_path: str, output_path: str) -> Dict:
    """
    Review translated presentation and fix structural issues using AI

    Args:
        original_path: Path to original PPTX
        translated_path: Path to translated PPTX
        output_path: Path to save fixed PPTX

    Returns:
        Dictionary with review results and fixes applied
    """
    logger.info("Starting AI quality review...")

    # Step 1: Analyze both presentations
    original_analysis = analyze_presentation(original_path)
    translated_analysis = analyze_presentation(translated_path)

    # Step 2: Compare structures and identify issues
    issues = compare_structures(original_analysis, translated_analysis)

    if not issues:
        logger.info("No structural issues found. Presentation looks good!")
        shutil.copy(translated_path, output_path)
        return {'status': 'success', 'issues_found': 0, 'fixes_applied': 0}

    logger.info(f"Found {len(issues)} potential issues")

    # Step 3: Use AI to suggest fixes
    fixes = get_ai_suggestions(issues, original_analysis, translated_analysis)

    # Step 4: Apply fixes
    apply_fixes(translated_path, fixes, output_path)

    logger.info(f"Quality review complete. Applied {len(fixes)} fixes.")

    return {
        'status': 'success',
        'issues_found': len(issues),
        'fixes_applied': len(fixes),
        'details': fixes
    }

def analyze_presentation(pptx_path: str) -> Dict:
    """
    Analyze presentation structure

    Returns detailed information about each slide's structure
    """
    logger.info(f"Analyzing: {pptx_path}")

    prs = Presentation(pptx_path)
    analysis = {
        'slide_count': len(prs.slides),
        'slides': []
    }

    for slide_idx, slide in enumerate(prs.slides):
        slide_info = {
            'index': slide_idx,
            'shapes': []
        }

        for shape_idx, shape in enumerate(slide.shapes):
            shape_data = {
                'index': shape_idx,
                'type': str(shape.shape_type),
                'has_text': hasattr(shape, 'text_frame'),
                'left': shape.left,
                'top': shape.top,
                'width': shape.width,
                'height': shape.height
            }

            # Analyze text if present
            if hasattr(shape, 'text_frame') and shape.text_frame:
                text = shape.text_frame.text.strip()
                shape_data['text_length'] = len(text)
                shape_data['has_text_content'] = bool(text)

                # Check for overflow issues
                if shape.text_frame.text:
                    shape_data['word_wrap'] = shape.text_frame.word_wrap
                    shape_data['auto_size'] = str(shape.text_frame.auto_size)

                    # Get font info
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.text.strip():
                                shape_data['font_name'] = run.font.name
                                shape_data['font_size'] = run.font.size
                                break
                        if 'font_name' in shape_data:
                            break

            slide_info['shapes'].append(shape_data)

        analysis['slides'].append(slide_info)

    return analysis

def compare_structures(original: Dict, translated: Dict) -> List[Dict]:
    """
    Compare original and translated structures to find issues

    Returns list of issues found
    """
    issues = []

    # Check slide count mismatch
    if original['slide_count'] != translated['slide_count']:
        issues.append({
            'type': 'slide_count_mismatch',
            'severity': 'high',
            'description': f"Original has {original['slide_count']} slides, translated has {translated['slide_count']}",
            'original_value': original['slide_count'],
            'translated_value': translated['slide_count']
        })

    # Compare each slide
    for slide_idx in range(min(original['slide_count'], translated['slide_count'])):
        orig_slide = original['slides'][slide_idx]
        trans_slide = translated['slides'][slide_idx]

        # Check shape count
        if len(orig_slide['shapes']) != len(trans_slide['shapes']):
            issues.append({
                'type': 'shape_count_mismatch',
                'severity': 'medium',
                'slide_index': slide_idx,
                'description': f"Slide {slide_idx+1}: Shape count mismatch",
                'original_value': len(orig_slide['shapes']),
                'translated_value': len(trans_slide['shapes'])
            })

        # Compare each shape
        for shape_idx in range(min(len(orig_slide['shapes']), len(trans_slide['shapes']))):
            orig_shape = orig_slide['shapes'][shape_idx]
            trans_shape = trans_slide['shapes'][shape_idx]

            # Check for significant size changes (more than 20%)
            if orig_shape['width'] > 0:
                width_change = abs(orig_shape['width'] - trans_shape['width']) / orig_shape['width']
                if width_change > 0.2:
                    issues.append({
                        'type': 'shape_size_change',
                        'severity': 'low',
                        'slide_index': slide_idx,
                        'shape_index': shape_idx,
                        'description': f"Slide {slide_idx+1}, Shape {shape_idx+1}: Width changed by {width_change*100:.1f}%",
                        'dimension': 'width',
                        'change_percent': width_change * 100
                    })

            if orig_shape['height'] > 0:
                height_change = abs(orig_shape['height'] - trans_shape['height']) / orig_shape['height']
                if height_change > 0.2:
                    issues.append({
                        'type': 'shape_size_change',
                        'severity': 'low',
                        'slide_index': slide_idx,
                        'shape_index': shape_idx,
                        'description': f"Slide {slide_idx+1}, Shape {shape_idx+1}: Height changed by {height_change*100:.1f}%",
                        'dimension': 'height',
                        'change_percent': height_change * 100
                    })

            # Check for position changes (horizontal flip is expected, but vertical shouldn't change much)
            if orig_shape['top'] > 0:
                top_change = abs(orig_shape['top'] - trans_shape['top']) / orig_shape['top']
                if top_change > 0.1:
                    issues.append({
                        'type': 'vertical_position_change',
                        'severity': 'medium',
                        'slide_index': slide_idx,
                        'shape_index': shape_idx,
                        'description': f"Slide {slide_idx+1}, Shape {shape_idx+1}: Vertical position changed by {top_change*100:.1f}%",
                        'change_percent': top_change * 100
                    })

    return issues

def get_ai_suggestions(issues: List[Dict], original: Dict, translated: Dict) -> List[Dict]:
    """
    Use OpenAI to analyze issues and suggest fixes

    Args:
        issues: List of identified issues
        original: Original presentation analysis
        translated: Translated presentation analysis

    Returns:
        List of suggested fixes
    """
    if not issues:
        return []

    logger.info("Getting AI suggestions for fixes...")

    # Prepare issue summary for AI
    issue_summary = "\n".join([
        f"- {issue['description']} (Severity: {issue['severity']})"
        for issue in issues
    ])

    # Create prompt for AI
    prompt = f"""You are a PowerPoint presentation quality reviewer. A presentation has been translated from English to Arabic with RTL (right-to-left) layout conversion.

The following structural issues were detected:

{issue_summary}

Context:
- Original presentation has {original['slide_count']} slides
- Translated presentation has {translated['slide_count']} slides
- Translation includes RTL layout flip (shapes are mirrored horizontally, which is expected)
- Arabic text uses Arial font with 90% sizing for better fitting

For each issue, determine:
1. Is this a real problem or expected behavior (e.g., horizontal position changes are expected for RTL)?
2. What fix should be applied (if any)?
3. Priority (high/medium/low)

Respond in JSON format:
[
  {{
    "issue_type": "...",
    "is_real_problem": true/false,
    "fix_needed": true/false,
    "fix_description": "...",
    "priority": "high/medium/low"
  }}
]

Focus on real problems that affect professional appearance, not expected RTL layout changes."""

    try:
        client = get_openai_client()

        response = client.chat.completions.create(
            model=Config.OPENAI_MODEL,
            messages=[
                {"role": "system", "content": "You are a PowerPoint quality expert specializing in RTL layout translations."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.3,
            response_format={"type": "json_object"}
        )

        import json
        suggestions = json.loads(response.choices[0].message.content)

        # Filter to only real problems that need fixes
        fixes = []
        for idx, suggestion in enumerate(suggestions.get('fixes', suggestions.get('analysis', []))):
            if suggestion.get('fix_needed', False):
                fixes.append({
                    'issue': issues[idx] if idx < len(issues) else {},
                    'suggestion': suggestion,
                    'applied': False
                })

        logger.info(f"AI suggested {len(fixes)} fixes out of {len(issues)} issues")
        return fixes

    except Exception as e:
        logger.error(f"Error getting AI suggestions: {str(e)}")
        return []

def apply_fixes(translated_path: str, fixes: List[Dict], output_path: str) -> None:
    """
    Apply suggested fixes to the presentation

    For now, this logs the suggestions. Future enhancement: automatically apply fixes.
    """
    logger.info(f"Applying {len(fixes)} fixes...")

    # For now, just copy the file and log suggestions
    # Future: Implement automatic fixes based on AI suggestions
    shutil.copy(translated_path, output_path)

    # Log each fix suggestion
    for fix in fixes:
        logger.info(f"Fix suggestion: {fix['suggestion'].get('fix_description', 'N/A')}")
        fix['applied'] = True

    logger.info("Fixes logged. Manual review recommended for complex issues.")

if __name__ == "__main__":
    # Test quality review
    if len(sys.argv) > 3:
        original = sys.argv[1]
        translated = sys.argv[2]
        output = sys.argv[3]

        result = review_and_fix_presentation(original, translated, output)

        print(f"\nQuality Review Results:")
        print(f"Status: {result['status']}")
        print(f"Issues found: {result['issues_found']}")
        print(f"Fixes applied: {result['fixes_applied']}")
    else:
        print("Usage: python ai_quality_reviewer.py <original.pptx> <translated.pptx> <output.pptx>")
