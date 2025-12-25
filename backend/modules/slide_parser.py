"""
Slide Parser Module
Extracts slide structure, text elements, and hierarchy from PowerPoint slides
"""
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from typing import Dict, List, Any
import sys
import os

# Add parent directory to path
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from utils.logger import setup_logger

logger = setup_logger(__name__)

def extract_slide_structure(pptx_path: str, slide_index: int = 0) -> Dict[str, Any]:
    """
    Extract structure and content from a PowerPoint slide

    Args:
        pptx_path: Path to PowerPoint file
        slide_index: Index of slide to extract (default: 0 = first slide)

    Returns:
        Dictionary containing:
        {
            "slide_index": 0,
            "slide_width": 9144000,
            "slide_height": 6858000,
            "elements": [
                {
                    "element_id": "shape_0",
                    "type": "title",
                    "text": "Our Strategic Approach",
                    "shape_id": 2,
                    "position": {"left": 100, "top": 50, "width": 800, "height": 60},
                    "level": 0
                },
                ...
            ]
        }
    """
    logger.info(f"Parsing PowerPoint file: {pptx_path}")

    try:
        # Load presentation
        prs = Presentation(pptx_path)

        # Check if slide index exists
        if slide_index >= len(prs.slides):
            raise ValueError(f"Slide index {slide_index} out of range. Total slides: {len(prs.slides)}")

        slide = prs.slides[slide_index]

        # Initialize structure
        structure = {
            "slide_index": slide_index,
            "slide_width": prs.slide_width,
            "slide_height": prs.slide_height,
            "elements": []
        }

        # Extract elements from shapes (including nested groups)
        element_counter = [0]  # Use list to allow modification in nested function

        def extract_from_shape(shape, parent_id=""):
            """Recursively extract text from shape (handles groups, tables, text)"""
            from pptx.enum.shapes import MSO_SHAPE_TYPE

            # Handle GROUPED SHAPES - recurse into nested shapes
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for sub_idx, sub_shape in enumerate(shape.shapes):
                    sub_id = f"{parent_id}_group{sub_idx}" if parent_id else f"group{sub_idx}"
                    extract_from_shape(sub_shape, sub_id)
                return

            # Handle TABLES
            if hasattr(shape, 'has_table') and shape.has_table:
                for row_idx, row in enumerate(shape.table.rows):
                    for col_idx, cell in enumerate(row.cells):
                        text = cell.text.strip()
                        if text:
                            elem_id = f"table_{element_counter[0]}"
                            element_counter[0] += 1
                            structure["elements"].append({
                                "element_id": elem_id,
                                "type": "table_cell",
                                "text": text,
                                "shape_id": getattr(shape, 'shape_id', 0),
                                "position": {
                                    "left": shape.left if shape.left is not None else 0,
                                    "top": shape.top if shape.top is not None else 0,
                                    "width": shape.width if shape.width is not None else 0,
                                    "height": shape.height if shape.height is not None else 0
                                },
                                "level": 0
                            })
                return

            # Handle TEXT FRAMES (regular text boxes)
            if not shape.has_text_frame:
                return

            # Determine element type (needs shape index)
            element_type = _determine_element_type(shape, element_counter[0])

            # Extract text content
            if element_type == "bullet_group":
                bullets = _extract_bullets(shape)
                if bullets:
                    # Only increment counter when actually adding element
                    shape_id = f"shape_{element_counter[0]}"
                    element_counter[0] += 1
                    structure["elements"].append({
                        "element_id": shape_id,
                        "type": element_type,
                        "bullets": bullets,
                        "shape_id": getattr(shape, 'shape_id', 0),
                        "position": {
                            "left": shape.left if shape.left is not None else 0,
                            "top": shape.top if shape.top is not None else 0,
                            "width": shape.width if shape.width is not None else 0,
                            "height": shape.height if shape.height is not None else 0
                        }
                    })
            else:
                # Single text element
                text = shape.text.strip()
                if text:
                    # Only increment counter when actually adding element
                    shape_id = f"shape_{element_counter[0]}"
                    element_counter[0] += 1
                    structure["elements"].append({
                        "element_id": shape_id,
                        "type": element_type,
                        "text": text,
                        "shape_id": getattr(shape, 'shape_id', 0),
                        "position": {
                            "left": shape.left if shape.left is not None else 0,
                            "top": shape.top if shape.top is not None else 0,
                            "width": shape.width if shape.width is not None else 0,
                            "height": shape.height if shape.height is not None else 0
                        },
                        "level": 0
                    })

        # Process all top-level shapes
        for shape in slide.shapes:
            extract_from_shape(shape)

        logger.info(f"Extracted {len(structure['elements'])} elements from slide")
        return structure

    except Exception as e:
        logger.error(f"Error parsing slide: {str(e)}", exc_info=True)
        raise

def _determine_element_type(shape, shape_idx: int) -> str:
    """
    Determine the type of element (title, header, bullet_group, text_box)

    Args:
        shape: PowerPoint shape object
        shape_idx: Index of shape on slide

    Returns:
        Element type string
    """
    # Check if it's a title placeholder
    if shape.is_placeholder:
        ph_type = shape.placeholder_format.type
        if ph_type == 1:  # PP_PLACEHOLDER.TITLE
            return "title"
        elif ph_type == 2:  # PP_PLACEHOLDER.BODY
            return "bullet_group"

    # Check if shape has bullets
    if shape.has_text_frame:
        text_frame = shape.text_frame
        has_bullets = any(
            paragraph.level > 0 or (len(text_frame.paragraphs) > 1 and paragraph.text.strip())
            for paragraph in text_frame.paragraphs
        )
        if has_bullets or len(text_frame.paragraphs) > 2:
            return "bullet_group"

    # Check position to determine if it's a header (near top)
    # Skip if shape.top is None (can happen with certain shape types)
    if shape_idx <= 1 and shape.top is not None and shape.top < 2000000:  # Near top (measured in EMUs)
        return "header"

    # Default to text_box
    return "text_box"

def _extract_bullets(shape) -> List[Dict[str, Any]]:
    """
    Extract bullet points from a text frame

    Args:
        shape: PowerPoint shape with text_frame

    Returns:
        List of bullet dictionaries
    """
    bullets = []

    if not shape.has_text_frame:
        return bullets

    text_frame = shape.text_frame

    for para_idx, paragraph in enumerate(text_frame.paragraphs):
        text = paragraph.text.strip()
        if text:  # Only include non-empty paragraphs
            bullets.append({
                "text": text,
                "level": paragraph.level,
                "index": para_idx
            })

    return bullets

def get_text_by_element_id(structure: Dict[str, Any], element_id: str) -> str:
    """
    Get text content by element ID

    Args:
        structure: Slide structure from extract_slide_structure
        element_id: Element ID to retrieve

    Returns:
        Text content or empty string
    """
    for element in structure["elements"]:
        if element["element_id"] == element_id:
            if "text" in element:
                return element["text"]
            elif "bullets" in element:
                return "\n".join([b["text"] for b in element["bullets"]])

    return ""

if __name__ == "__main__":
    # Test the parser
    import sys
    if len(sys.argv) > 1:
        test_file = sys.argv[1]
        result = extract_slide_structure(test_file)
        import json
        print(json.dumps(result, indent=2))
    else:
        print("Usage: python slide_parser.py <pptx_file>")
