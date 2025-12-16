"""
PDF Converter Module
Converts PDF files to PowerPoint presentations for translation
"""
import os
import sys
import tempfile
from pathlib import Path
from typing import Optional

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from utils.logger import setup_logger

logger = setup_logger(__name__)

def convert_pdf_to_pptx(pdf_path: str, output_path: Optional[str] = None) -> str:
    """
    Convert PDF file to PowerPoint presentation

    Args:
        pdf_path: Path to input PDF file
        output_path: Path to save PPTX (optional, auto-generated if not provided)

    Returns:
        Path to converted PPTX file
    """
    logger.info(f"Converting PDF to PPTX: {pdf_path}")

    # Generate output path if not provided
    if output_path is None:
        pdf_name = Path(pdf_path).stem
        output_path = str(Path(pdf_path).parent / f"{pdf_name}_converted.pptx")

    try:
        # Try using pdf2image + pptx approach
        return _convert_pdf_via_images(pdf_path, output_path)

    except Exception as e:
        logger.error(f"PDF conversion failed: {str(e)}")
        raise Exception(f"Failed to convert PDF to PowerPoint: {str(e)}")

def _convert_pdf_via_images(pdf_path: str, output_path: str) -> str:
    """
    Convert PDF to PPTX by converting pages to images and inserting into slides

    Args:
        pdf_path: Path to PDF file
        output_path: Path to save PPTX

    Returns:
        Path to converted PPTX
    """
    try:
        import fitz  # PyMuPDF
        from pptx import Presentation
        from pptx.util import Inches
        from PIL import Image
        import io

        logger.info("Converting PDF pages to images...")

        # Open PDF with PyMuPDF
        pdf_document = fitz.open(pdf_path)
        page_count = len(pdf_document)

        logger.info(f"Extracted {page_count} pages from PDF")

        # Convert pages to images
        images = []
        for page_num in range(page_count):
            page = pdf_document[page_num]
            # Render page to image (300 DPI)
            pix = page.get_pixmap(matrix=fitz.Matrix(300/72, 300/72))
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            images.append(img)

        pdf_document.close()

        logger.info(f"Converted {len(images)} pages to images")

        # Create PowerPoint presentation
        prs = Presentation()

        # Set slide dimensions (16:9 aspect ratio)
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        temp_dir = tempfile.mkdtemp()

        # Reopen PDF to extract text
        pdf_document = fitz.open(pdf_path)

        for idx, image in enumerate(images):
            logger.info(f"Processing page {idx + 1}/{len(images)}...")

            # Save image temporarily
            img_path = os.path.join(temp_dir, f'page_{idx}.png')
            image.save(img_path, 'PNG')

            # Extract text from PDF page
            try:
                page = pdf_document[idx]
                text = page.get_text()
                logger.info(f"  Extracted {len(text)} characters from page")
            except Exception as e:
                logger.warning(f"  Text extraction failed: {str(e)}")
                text = ""

            # Add blank slide
            blank_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(blank_layout)

            # Add image as background
            left = Inches(0)
            top = Inches(0)
            width = prs.slide_width
            height = prs.slide_height

            pic = slide.shapes.add_picture(img_path, left, top, width=width, height=height)

            # Move picture to background
            slide.shapes._spTree.remove(pic._element)
            slide.shapes._spTree.insert(2, pic._element)

            # Add text boxes with extracted text if available
            if text.strip():
                # Add text box
                left = Inches(0.5)
                top = Inches(0.5)
                width = Inches(9)
                height = Inches(6.5)

                textbox = slide.shapes.add_textbox(left, top, width, height)
                text_frame = textbox.text_frame
                text_frame.text = text.strip()

                # Make text box transparent
                textbox.fill.background()

        # Save presentation
        prs.save(output_path)

        # Close PDF document
        pdf_document.close()

        # Cleanup
        import shutil
        shutil.rmtree(temp_dir)

        logger.info(f"PDF converted successfully: {output_path}")
        return output_path

    except ImportError as e:
        logger.error(f"Missing required library: {str(e)}")
        raise Exception(
            "PDF conversion requires additional libraries. "
            "Install with: pip install pymupdf python-pptx pillow"
        )

def is_pdf_file(filename: str) -> bool:
    """
    Check if file is a PDF

    Args:
        filename: Name of the file

    Returns:
        True if PDF file
    """
    return filename.lower().endswith('.pdf')

if __name__ == "__main__":
    # Test PDF conversion
    if len(sys.argv) > 1:
        pdf_file = sys.argv[1]
        output_file = sys.argv[2] if len(sys.argv) > 2 else None

        result = convert_pdf_to_pptx(pdf_file, output_file)
        print(f"Conversion complete: {result}")
    else:
        print("Usage: python pdf_converter.py <input.pdf> [output.pptx]")
