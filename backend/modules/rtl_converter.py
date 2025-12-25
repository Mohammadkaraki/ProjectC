"""
RTL Converter Module
Converts PowerPoint slide layout from Left-to-Right (LTR) to Right-to-Left (RTL)
This is the most complex module as it requires XML manipulation
"""
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE_TYPE
from lxml import etree
from typing import Any
import sys
import os

# Add parent directory to path
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from config import Config
from utils.logger import setup_logger

logger = setup_logger(__name__)

# XML namespaces for PowerPoint
NAMESPACES = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
}

def flip_to_rtl_layout(input_path: str, output_path: str) -> None:
    """
    Convert PowerPoint slide from LTR to RTL layout

    This function prepares the slide for RTL content by:
    - Setting paragraph-level RTL direction (for proper text rendering)
    - Text alignment (RIGHT)
    - Shape positions (mirrored)
    - Reading order

    NOTE: We do NOT set slide-level RTL as it can cause character reversal
    with Arabic text that's already in correct logical order from AI translation.

    IMPORTANT: Chart elements should be grouped BEFORE calling this function.
    Grouped chart elements will flip as a single unit, preserving their internal layout.

    Args:
        input_path: Path to input .pptx file
        output_path: Path to save RTL-converted .pptx file
    """
    logger.info(f"Converting slide to RTL layout: {input_path}")

    try:
        # Load presentation
        prs = Presentation(input_path)

        # Process all slides with normal RTL conversion
        for slide_idx, slide in enumerate(prs.slides):
            logger.info(f"Processing slide {slide_idx + 1}/{len(prs.slides)}")

            # Flip shape positions horizontally (mirror the slide)
            chart_regions = _get_chart_regions(slide)
            slide_width = prs.slide_width
            for shape in slide.shapes:
                _flip_shape_position(shape, slide_width, chart_regions)

            # Set text alignment to RIGHT and paragraph-level RTL
            for shape in slide.shapes:
                if shape.has_text_frame:
                    _set_text_rtl_and_alignment(shape)

        # Save modified presentation
        prs.save(output_path)
        logger.info(f"RTL conversion complete. Saved to: {output_path}")

    except Exception as e:
        logger.error(f"Error during RTL conversion: {str(e)}", exc_info=True)
        raise

def _get_shape_id(shape) -> str:
    """
    Get a unique identifier for a shape based on its PowerPoint ID

    Args:
        shape: PowerPoint shape object

    Returns:
        Unique string identifier for the shape
    """
    try:
        # Try to get the shape's PowerPoint ID from XML
        nvSpPr = shape._element.find('.//p:cNvPr', namespaces=NAMESPACES)
        if nvSpPr is not None:
            shape_id = nvSpPr.get('id')
            shape_name = nvSpPr.get('name', '')
            return f"{shape_id}_{shape_name}"
    except:
        pass

    # Fallback to object id (memory address)
    return str(id(shape._element))

def group_chart_elements(slide) -> int:
    """
    Group chart and all related elements (text boxes, labels, etc.) on a slide

    This function:
    1. Finds all charts on the slide
    2. For each chart, finds related elements (nearby shapes, text boxes)
    3. Groups the chart + related elements together

    Args:
        slide: PowerPoint slide object

    Returns:
        Number of groups created
    """
    charts_found = []

    # Step 1: Find all charts on the slide
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.CHART:
            charts_found.append(shape)

    if not charts_found:
        return 0

    logger.info(f"Found {len(charts_found)} chart(s) on slide - grouping related elements")

    groups_created = 0
    already_grouped = set()  # Track shapes that have been grouped

    # Step 2: Build a map of all shapes to their nearest chart
    shape_to_nearest_chart = {}
    # STRICT threshold: Only group elements that are EXACTLY touching/overlapping the chart (distance = 0)
    # This prevents grouping nearby text that's not actually part of the chart
    proximity_threshold = 0  # 0 inches - only elements that are touching/overlapping (distance = 0.0000)

    for shape in slide.shapes:
        shape_id = _get_shape_id(shape)

        # Skip charts themselves
        if shape.shape_type == MSO_SHAPE_TYPE.CHART:
            continue

        # Skip pre-existing groups (they may contain unrelated elements)
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            continue

        # Skip title placeholders (slide titles should never be grouped with charts)
        if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
            if hasattr(shape, 'name') and 'title' in shape.name.lower():
                continue

        # Find nearest chart for this shape
        min_distance = float('inf')
        nearest_chart_idx = None

        for chart_idx, chart_shape in enumerate(charts_found):
            distance = _distance_to_chart(shape, chart_shape)
            if distance < min_distance:
                min_distance = distance
                nearest_chart_idx = chart_idx

        # Only associate with chart if within threshold
        if min_distance <= proximity_threshold:
            shape_to_nearest_chart[shape_id] = (nearest_chart_idx, min_distance)
            shape_name = shape.name if hasattr(shape, 'name') else 'unnamed'
            logger.info(f"  Shape '{shape_name}' [ID:{shape_id}] -> Chart {nearest_chart_idx + 1} (distance: {min_distance/914400:.2f} in) [ADDED TO MAP]")
        else:
            if hasattr(shape, 'text') and shape.text and len(shape.text) > 0:
                shape_name = shape.name if hasattr(shape, 'name') else 'unnamed'
                logger.info(f"  Shape '{shape_name}' text='{shape.text[:30]}' -> TOO FAR (distance: {min_distance/914400:.2f} in > threshold {proximity_threshold/914400:.2f} in) [NOT ADDED]")

    # Step 3: For each chart, group it with its associated elements
    for chart_idx, chart_shape in enumerate(charts_found):
        # Skip if this chart was already grouped
        chart_id = _get_shape_id(chart_shape)
        if chart_id in already_grouped:
            logger.info(f"  Chart {chart_idx + 1} already grouped, skipping")
            continue

        # Find related elements (assigned to this chart)
        related_shapes = []
        related_shapes.append(chart_shape)  # Start with the chart itself
        already_grouped.add(chart_id)

        # Add shapes that belong to this chart
        for shape in slide.shapes:
            shape_id = _get_shape_id(shape)

            # Skip if already grouped
            if shape_id in already_grouped:
                continue

            # Skip groups (defensive check - they should already be filtered out)
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                continue

            # Skip charts (defensive check)
            if shape.shape_type == MSO_SHAPE_TYPE.CHART:
                continue

            # Skip title placeholders (defensive check)
            if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                if hasattr(shape, 'name') and 'title' in shape.name.lower():
                    continue

            # Check if this shape is assigned to the current chart
            if shape_id in shape_to_nearest_chart:
                assigned_chart_idx, distance = shape_to_nearest_chart[shape_id]
                if assigned_chart_idx == chart_idx:
                    shape_name = shape.name if hasattr(shape, 'name') else 'unnamed'
                    logger.info(f"  [ADDING] '{shape_name}' [ID:{shape_id}] to chart {chart_idx + 1} (was in map with distance {distance/914400:.2f} in)")
                    related_shapes.append(shape)
                    already_grouped.add(shape_id)

        # Group the chart and related shapes
        if len(related_shapes) > 1:  # Only group if there are multiple elements
            try:
                # Log what we're about to group
                logger.info(f"  Grouping chart {chart_idx + 1} with {len(related_shapes)} total shapes:")
                for rs in related_shapes:
                    shape_name = rs.name if hasattr(rs, 'name') else 'unnamed'
                    shape_type_name = str(rs.shape_type)
                    logger.info(f"    - {shape_name} (type: {shape_type_name})")

                # Get the shapes collection
                shapes_collection = slide.shapes

                # Group the shapes (chart + related elements)
                # Note: python-pptx doesn't have a direct group method, we need to do it via XML
                _group_shapes_xml(slide, related_shapes)
                groups_created += 1
                logger.info(f"  Successfully grouped chart {chart_idx + 1}")
            except Exception as e:
                logger.warning(f"  Could not group chart elements: {str(e)}")
        else:
            logger.info(f"  Chart {chart_idx + 1} has no related elements to group")

    return groups_created

def _distance_to_chart(shape, chart) -> float:
    """
    Calculate the minimum distance from a shape to a chart

    Args:
        shape: PowerPoint shape
        chart: PowerPoint chart shape

    Returns:
        Minimum distance in EMUs
    """
    if shape.left is None or shape.top is None or chart.left is None or chart.top is None:
        return float('inf')

    # Get shape boundaries
    shape_left = shape.left
    shape_top = shape.top
    shape_right = shape.left + (shape.width if shape.width else 0)
    shape_bottom = shape.top + (shape.height if shape.height else 0)

    # Get chart boundaries
    chart_left = chart.left
    chart_top = chart.top
    chart_right = chart.left + (chart.width if chart.width else 0)
    chart_bottom = chart.top + (chart.height if chart.height else 0)

    # Calculate minimum distance between rectangles
    # If they overlap, distance is 0
    if (shape_right >= chart_left and shape_left <= chart_right and
        shape_bottom >= chart_top and shape_top <= chart_bottom):
        return 0

    # Calculate distance to nearest edge
    dx = 0
    if shape_right < chart_left:
        dx = chart_left - shape_right
    elif shape_left > chart_right:
        dx = shape_left - chart_right

    dy = 0
    if shape_bottom < chart_top:
        dy = chart_top - shape_bottom
    elif shape_top > chart_bottom:
        dy = shape_top - chart_bottom

    # Return Euclidean distance
    return (dx * dx + dy * dy) ** 0.5

def _is_shape_near_region(shape, region: dict, proximity: int = 700000) -> bool:
    """
    Check if a shape is near a region

    Args:
        shape: PowerPoint shape
        region: Dictionary with 'left', 'top', 'right', 'bottom' keys
        proximity: Distance threshold in EMUs

    Returns:
        True if shape is near the region
    """
    if shape.left is None or shape.top is None:
        return False

    shape_left = shape.left
    shape_top = shape.top

    # Check if shape is within proximity of region
    if (shape_left >= region['left'] - proximity and
        shape_left <= region['right'] + proximity and
        shape_top >= region['top'] - proximity and
        shape_top <= region['bottom'] + proximity):
        return True

    return False

def _group_shapes_xml(slide, shapes_to_group: list) -> None:
    """
    Group shapes together using XML manipulation

    Creates a group shape (grpSp) and moves the specified shapes into it.

    Args:
        slide: PowerPoint slide object
        shapes_to_group: List of shapes to group
    """
    try:
        if len(shapes_to_group) < 2:
            logger.warning("Cannot group less than 2 shapes")
            return

        # Get slide's shape tree (spTree)
        slide_element = slide._element
        spTree = slide_element.find('.//p:spTree', namespaces=NAMESPACES)

        if spTree is None:
            logger.error("Could not find shape tree in slide")
            return

        # Calculate bounding box for the group
        min_left = min(s.left for s in shapes_to_group if s.left is not None)
        min_top = min(s.top for s in shapes_to_group if s.top is not None)
        max_right = max(s.left + s.width for s in shapes_to_group if s.left is not None and s.width is not None)
        max_bottom = max((s.top if s.top is not None else 0) + (s.height if s.height is not None else 0) for s in shapes_to_group)

        group_width = max_right - min_left
        group_height = max_bottom - min_top

        # Create group shape element (grpSp)
        grpSp = etree.Element(f"{{{NAMESPACES['p']}}}grpSp")

        # Add nvGrpSpPr (non-visual group shape properties)
        nvGrpSpPr = etree.SubElement(grpSp, f"{{{NAMESPACES['p']}}}nvGrpSpPr")
        cNvPr = etree.SubElement(nvGrpSpPr, f"{{{NAMESPACES['p']}}}cNvPr")
        cNvPr.set('id', str(len(spTree) + 1))
        cNvPr.set('name', f'Group {len(spTree) + 1}')
        etree.SubElement(nvGrpSpPr, f"{{{NAMESPACES['p']}}}cNvGrpSpPr")
        etree.SubElement(nvGrpSpPr, f"{{{NAMESPACES['p']}}}nvPr")

        # Add grpSpPr (group shape properties) with transform
        grpSpPr = etree.SubElement(grpSp, f"{{{NAMESPACES['p']}}}grpSpPr")
        xfrm = etree.SubElement(grpSpPr, f"{{{NAMESPACES['a']}}}xfrm")

        # Offset (group position)
        off = etree.SubElement(xfrm, f"{{{NAMESPACES['a']}}}off")
        off.set('x', str(min_left))
        off.set('y', str(min_top))

        # Extents (group size)
        ext = etree.SubElement(xfrm, f"{{{NAMESPACES['a']}}}ext")
        ext.set('cx', str(group_width))
        ext.set('cy', str(group_height))

        # Child offset (always 0,0 for groups)
        chOff = etree.SubElement(xfrm, f"{{{NAMESPACES['a']}}}chOff")
        chOff.set('x', '0')
        chOff.set('y', '0')

        # Child extents (same as group extents)
        chExt = etree.SubElement(xfrm, f"{{{NAMESPACES['a']}}}chExt")
        chExt.set('cx', str(group_width))
        chExt.set('cy', str(group_height))

        # FIRST: Capture all absolute positions BEFORE any XML manipulation
        shape_positions = []
        for shape in shapes_to_group:
            shape_positions.append({
                'element': shape._element,
                'left': shape.left if shape.left is not None else 0,
                'top': shape.top if shape.top is not None else 0
            })

        # THEN: Move shapes into the group and adjust their coordinates
        for shape_info in shape_positions:
            shape_element = shape_info['element']
            abs_left = shape_info['left']
            abs_top = shape_info['top']

            # Calculate relative position to group
            rel_left = abs_left - min_left
            rel_top = abs_top - min_top

            logger.debug(f"  Shape absolute: ({abs_left}, {abs_top})")
            logger.debug(f"  Group min: ({min_left}, {min_top})")
            logger.debug(f"  Calculated relative: ({rel_left}, {rel_top})")

            # Adjust shape position in XML to be relative to group
            # Need to check different XML structures for different shape types
            shape_spPr = shape_element.find('.//p:spPr', namespaces=NAMESPACES)

            # For charts, the structure might be different - check graphicFrame
            if shape_spPr is None:
                # Try graphicFrame for charts
                shape_spPr = shape_element.find('.//p:xfrm', namespaces=NAMESPACES)
                if shape_spPr is not None:
                    off = shape_spPr.find('.//a:off', namespaces=NAMESPACES)
                    if off is not None:
                        logger.debug(f"  Found chart xfrm/off - setting to ({rel_left}, {rel_top})")
                        off.set('x', str(int(rel_left)))
                        off.set('y', str(int(rel_top)))
            else:
                shape_xfrm = shape_spPr.find('.//a:xfrm', namespaces=NAMESPACES)
                if shape_xfrm is not None:
                    shape_off = shape_xfrm.find('.//a:off', namespaces=NAMESPACES)
                    if shape_off is not None:
                        logger.debug(f"  Found shape spPr/xfrm/off - setting to ({rel_left}, {rel_top})")
                        # Set relative coordinates
                        shape_off.set('x', str(int(rel_left)))
                        shape_off.set('y', str(int(rel_top)))

            # Add shape to group
            grpSp.append(shape_element)

        # Remove shapes from their original position (they're now in the group)
        # (Already moved via append, XML automatically removes from old parent)

        # Add group to slide's shape tree
        spTree.append(grpSp)

        logger.info(f"  Successfully created group with {len(shapes_to_group)} shapes")

    except Exception as e:
        logger.error(f"Error in XML grouping: {str(e)}", exc_info=True)
        raise

def detect_chart_slides(pptx_path: str) -> list:
    """
    Detect which slides contain charts

    Args:
        pptx_path: Path to PowerPoint file

    Returns:
        List of slide indices (0-based) that contain charts
    """
    chart_slide_indices = []

    try:
        prs = Presentation(pptx_path)

        for slide_idx, slide in enumerate(prs.slides):
            has_chart = False
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.CHART:
                    has_chart = True
                    break

            if has_chart:
                chart_slide_indices.append(slide_idx)
                logger.info(f"Slide {slide_idx + 1} contains chart(s)")

        logger.info(f"Found {len(chart_slide_indices)} slide(s) with charts")

    except Exception as e:
        logger.error(f"Error detecting chart slides: {str(e)}")

    return chart_slide_indices

def _get_chart_regions(slide) -> list:
    """
    Get bounding regions of all charts on the slide

    Returns list of dictionaries with chart boundaries:
    [{'left': x, 'top': y, 'right': x2, 'bottom': y2}, ...]
    """
    chart_regions = []

    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.CHART:
            if shape.left is not None and shape.width is not None:
                chart_regions.append({
                    'left': shape.left,
                    'top': shape.top if shape.top is not None else 0,
                    'right': shape.left + shape.width,
                    'bottom': (shape.top if shape.top is not None else 0) + (shape.height if shape.height is not None else 0)
                })

    logger.info(f"Found {len(chart_regions)} chart(s) on slide")
    return chart_regions

def _adjust_overlapping_groups(slide, slide_width: int) -> None:
    """
    Detect and fix overlapping groups after RTL flip

    Only applies to slides with 2 or more groups (which typically contain charts).
    After RTL flipping, groups may overlap due to their original positions.
    This function detects overlaps and adds spacing between groups.

    Args:
        slide: PowerPoint slide object
        slide_width: Width of the slide in EMUs
    """
    # Collect all groups on the slide
    groups = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            groups.append(shape)

    # Only process slides with 2+ groups
    if len(groups) < 2:
        logger.debug(f"Slide has {len(groups)} group(s) - skipping overlap adjustment")
        return

    logger.info(f"Checking for overlaps among {len(groups)} groups on slide")

    # Sort groups by left position (left to right)
    groups_sorted = sorted(groups, key=lambda g: g.left if g.left else 0)

    # Minimum spacing between groups (0.2 inches)
    min_spacing = 182880  # 0.2 inches in EMUs

    # Check for overlaps and adjust positions
    # Only adjust if groups overlap BOTH horizontally AND vertically
    adjusted_count = 0
    for i in range(len(groups_sorted) - 1):
        current = groups_sorted[i]
        next_group = groups_sorted[i + 1]

        # Get horizontal boundaries
        current_left = current.left if current.left else 0
        current_width = current.width if current.width else 0
        current_right = current_left + current_width

        next_left = next_group.left if next_group.left else 0
        next_width = next_group.width if next_group.width else 0

        # Get vertical boundaries
        current_top = current.top if current.top else 0
        current_height = current.height if current.height else 0
        current_bottom = current_top + current_height

        next_top = next_group.top if next_group.top else 0
        next_height = next_group.height if next_group.height else 0
        next_bottom = next_top + next_height

        # Check if they overlap or are too close HORIZONTALLY
        horizontal_gap = next_left - current_right
        horizontal_overlap = horizontal_gap < min_spacing

        # Check if they overlap VERTICALLY (same vertical region)
        # Two groups overlap vertically if one's top is between the other's top and bottom
        vertical_overlap = not (current_bottom <= next_top or next_bottom <= current_top)

        # Only adjust if BOTH horizontal and vertical overlap exist
        if horizontal_overlap and vertical_overlap:
            # Calculate how much to shift the next group
            shift_amount = min_spacing - horizontal_gap
            new_left = next_left + shift_amount

            # Make sure we don't go off the slide
            if new_left + next_width <= slide_width:
                next_group.left = int(new_left)
                adjusted_count += 1
                logger.info(f"  Adjusted group '{next_group.name if hasattr(next_group, 'name') else 'unnamed'}' - added {shift_amount/914400:.2f}\" spacing")
            else:
                # If we can't shift right, try shifting the current group left
                shift_current = min_spacing - horizontal_gap
                new_current_left = current_left - shift_current

                if new_current_left >= 0:
                    current.left = int(new_current_left)
                    adjusted_count += 1
                    logger.info(f"  Adjusted group '{current.name if hasattr(current, 'name') else 'unnamed'}' leftward - added {shift_current/914400:.2f}\" spacing")
                else:
                    logger.warning(f"  Cannot adjust groups '{current.name if hasattr(current, 'name') else 'unnamed'}' and '{next_group.name if hasattr(next_group, 'name') else 'unnamed'}' - insufficient space")

    if adjusted_count > 0:
        logger.info(f"Auto-spacing: Adjusted {adjusted_count} group position(s) to prevent overlap")
    else:
        logger.info(f"Auto-spacing: No overlaps detected")

def _is_near_chart(shape, chart_regions, proximity=700000) -> bool:
    """
    Check if a shape is near any chart

    Args:
        shape: Shape to check
        chart_regions: List of chart boundary dictionaries
        proximity: Distance threshold in EMUs (default: 700000 = ~0.75 inches)
                   Captures chart titles, legends, and labels directly attached to the chart

    Returns:
        True if shape is near a chart, False otherwise
    """
    if not chart_regions:
        return False

    shape_left = shape.left if shape.left is not None else 0
    shape_top = shape.top if shape.top is not None else 0

    for chart in chart_regions:
        # Check if shape is within proximity of chart region
        # Use a large exclusion zone to capture all chart-related text, labels, legends, titles
        if (shape_left >= chart['left'] - proximity and
            shape_left <= chart['right'] + proximity and
            shape_top >= chart['top'] - proximity and
            shape_top <= chart['bottom'] + proximity):
            return True

    return False

def _set_slide_rtl_property(slide) -> None:
    """
    Set PowerPoint's built-in RTL property at the SLIDE level

    This is the correct way to enable RTL layout as per PowerPoint's design.
    Setting rtl="1" on the slide's p:cSld element activates PowerPoint's
    native right-to-left layout mode, which automatically:
    - Flips the reading order
    - Mirrors the slide layout
    - Sets text direction to RTL

    Args:
        slide: PowerPoint slide object
    """
    try:
        # Access slide XML element
        slide_element = slide._element

        # Find p:cSld (common slide data) element
        cSld = slide_element.find('.//p:cSld', namespaces=NAMESPACES)

        if cSld is not None:
            # Set RTL attribute on the slide (PowerPoint's built-in property)
            cSld.set('rtl', '1')  # 1 = RTL, 0 = LTR
            logger.info("✓ Set PowerPoint's built-in RTL property at slide level")
        else:
            logger.warning("Could not find p:cSld element to set RTL property")

    except Exception as e:
        logger.error(f"Error setting slide RTL property: {str(e)}")
        raise

def _set_text_rtl_and_alignment(shape) -> None:
    """
    Set text alignment to RIGHT and REMOVE any RTL attributes

    This sets:
    - RIGHT alignment (visual alignment to right side)
    - REMOVES rtl="1" attributes (they cause text reversal)

    NOTE: We REMOVE rtl="1" because Arabic text from LLM translation
    is already in correct logical order. The rtl="1" attribute would cause the text
    to be reversed again when displayed, making it unreadable.

    Args:
        shape: PowerPoint shape with text_frame
    """
    if not shape.has_text_frame:
        return

    text_frame = shape.text_frame

    for paragraph in text_frame.paragraphs:
        try:
            # Set RIGHT alignment
            paragraph.alignment = PP_ALIGN.RIGHT

            # REMOVE any existing RTL attribute
            _remove_rtl_from_paragraph(paragraph)
        except Exception as e:
            logger.warning(f"Could not set alignment for paragraph: {str(e)}")

def _set_text_alignment_right(shape) -> None:
    """
    Set text alignment to RIGHT for all paragraphs in a shape

    This complements the slide-level RTL property by ensuring
    text is visually aligned to the right side.

    Args:
        shape: PowerPoint shape with text_frame
    """
    if not shape.has_text_frame:
        return

    text_frame = shape.text_frame

    for paragraph in text_frame.paragraphs:
        try:
            paragraph.alignment = PP_ALIGN.RIGHT
        except Exception as e:
            logger.warning(f"Could not set alignment for paragraph: {str(e)}")

def _set_rtl_text_direction(shape) -> None:
    """
    Set text direction to RTL and align text to the right

    This function modifies:
    - Paragraph alignment (RIGHT)
    - Paragraph RTL property via XML (rtl="1")

    Args:
        shape: PowerPoint shape with text_frame
    """
    if not shape.has_text_frame:
        return

    text_frame = shape.text_frame

    for paragraph in text_frame.paragraphs:
        # Set alignment to RIGHT for RTL appearance
        paragraph.alignment = PP_ALIGN.RIGHT

        # Set RTL direction via XML manipulation
        try:
            _set_rtl_via_xml(paragraph)
        except Exception as e:
            logger.warning(f"Could not set RTL via XML: {str(e)}")

def _remove_rtl_from_paragraph(paragraph) -> None:
    """
    REMOVE RTL property from paragraph XML

    This removes the rtl="1" attribute that causes text reversal.
    Arabic text from LLM is already in correct logical order, so rtl="1"
    would reverse it and make it unreadable.

    Args:
        paragraph: PowerPoint paragraph object
    """
    try:
        # Get paragraph XML element
        p_element = paragraph._element

        # Get <a:pPr> (paragraph properties)
        pPr = p_element.find('.//a:pPr', namespaces=NAMESPACES)

        if pPr is not None and 'rtl' in pPr.attrib:
            # Remove RTL attribute
            del pPr.attrib['rtl']
            logger.debug("Removed RTL attribute from paragraph")

    except Exception as e:
        logger.warning(f"Could not remove RTL attribute: {str(e)}")

def _set_rtl_via_xml(paragraph) -> None:
    """
    Set RTL property in paragraph XML

    PowerPoint stores RTL information in the paragraph properties (pPr) XML element.
    We need to add rtl="1" attribute to the <a:pPr> element.

    Args:
        paragraph: PowerPoint paragraph object
    """
    # Get paragraph XML element
    p_element = paragraph._element

    # Get or create <a:pPr> (paragraph properties)
    pPr = p_element.find('.//a:pPr', namespaces=NAMESPACES)

    if pPr is None:
        # Create <a:pPr> if it doesn't exist
        # Insert before <a:r> (run) elements
        pPr = etree.Element(f"{{{NAMESPACES['a']}}}pPr")
        # Insert as first child
        p_element.insert(0, pPr)

    # Set RTL attribute
    pPr.set('rtl', '1')  # 1 = RTL, 0 = LTR

    logger.debug(f"Set RTL property via XML for paragraph")

def _flip_shape_position(shape, slide_width: int, chart_regions: list = None) -> None:
    """
    Flip text box position horizontally (mirror across slide center)
    PRESERVES width and height by explicitly setting them

    IMPORTANT: Charts and text boxes near charts are NOT flipped because:
    1. Charts are complex embedded objects with internal layout
    2. Text boxes near charts are likely labels/legends that belong to the chart
    3. Flipping them would misalign data labels from chart elements

    Calculation:
    new_left = slide_width - (old_left + width)

    Example:
    - Slide width: 9144000 EMUs (10 inches)
    - Shape at left=1000000, width=3000000
    - New position: 9144000 - (1000000 + 3000000) = 5144000

    Args:
        shape: PowerPoint shape to flip
        slide_width: Width of the slide in EMUs (English Metric Units)
        chart_regions: List of chart boundary regions (optional)
    """
    try:
        # SKIP CHARTS - they have internal layout that shouldn't be flipped
        if shape.shape_type == MSO_SHAPE_TYPE.CHART:
            logger.debug(f"Skipping chart shape - charts maintain original position for RTL")
            return

        # SKIP TEXT BOXES NEAR CHARTS - they likely belong to the chart
        if chart_regions and shape.has_text_frame and _is_near_chart(shape, chart_regions):
            logger.debug(f"Skipping text box near chart: '{shape.text.strip()[:30] if shape.text else 'empty'}'")
            return

        old_left = shape.left
        old_top = shape.top
        shape_width = shape.width
        shape_height = shape.height

        # Skip if position or size is None (can happen with certain shape types)
        if old_left is None or shape_width is None:
            logger.debug(f"Skipping shape with None position/size")
            return

        # Calculate new left position (mirror)
        new_left = slide_width - (old_left + shape_width)

        # Clamp to valid range (prevent negative positions)
        new_left = max(0, new_left)
        new_left = min(new_left, slide_width - shape_width)

        # Apply new position AND explicitly preserve width/height
        # This is CRITICAL for placeholders that don't have explicit dimensions
        shape.left = new_left
        shape.top = old_top if old_top is not None else 0
        shape.width = shape_width
        shape.height = shape_height if shape_height is not None else shape_width

        # MIRROR ARROWS: Flip arrow shapes horizontally for RTL
        # Check if this is an arrow shape (AutoShape with "Arrow" in name)
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and hasattr(shape, 'name'):
            if 'Arrow' in shape.name or 'arrow' in shape.name:
                _flip_shape_horizontally(shape)
                logger.debug(f"Mirrored arrow shape: {shape.name}")

        logger.debug(f"Flipped shape position: {old_left} → {new_left} (preserved w={shape_width}, h={shape_height})")

    except Exception as e:
        logger.warning(f"Could not flip shape position: {str(e)}")

def _flip_shape_in_place(shape) -> None:
    """
    Flip a shape horizontally without changing its position (for chart slides)

    This applies horizontal flip (mirror) to text boxes and shapes but keeps
    them in their exact position. Charts are NOT flipped.

    Uses PowerPoint's flipH XML attribute to mirror the shape.

    Args:
        shape: PowerPoint shape to flip in place
    """
    try:
        # SKIP CHARTS - they should not be flipped at all
        if shape.shape_type == MSO_SHAPE_TYPE.CHART:
            logger.debug(f"Skipping chart - charts are not flipped on chart slides")
            return

        # Apply horizontal flip via XML
        shape_element = shape._element
        spPr = shape_element.find('.//p:spPr', namespaces=NAMESPACES)

        if spPr is None:
            logger.debug(f"No spPr element found for shape")
            return

        # Find or create xfrm (transform) element
        xfrm = spPr.find('.//a:xfrm', namespaces=NAMESPACES)

        if xfrm is None:
            # Create xfrm element if it doesn't exist
            xfrm = etree.SubElement(spPr, f"{{{NAMESPACES['a']}}}xfrm")

        # Set flipH attribute to 1 (true)
        xfrm.set('flipH', '1')

        logger.debug(f"Applied horizontal flip to shape (kept in position)")

    except Exception as e:
        logger.warning(f"Could not flip shape in place: {str(e)}")

def _flip_shape_horizontally(shape) -> None:
    """
    Flip a shape horizontally for RTL layout (for arrows)

    For rotated arrows, we adjust the rotation angle instead of using flipH,
    because flipH on a rotated arrow produces incorrect results.

    Formula: new_rotation = 180° - old_rotation (mod 360°)

    Examples:
    - 0° (RIGHT →) becomes 180° (LEFT ←)
    - 315° (UP-RIGHT ↗) becomes 225° (UP-LEFT ↖)
    - 45° (DOWN-RIGHT ↘) becomes 135° (DOWN-LEFT ↙)

    Args:
        shape: PowerPoint shape to flip
    """
    try:
        # Get current rotation
        current_rotation = shape.rotation

        # Calculate mirrored rotation: 180° - current_rotation
        new_rotation = (180 - current_rotation) % 360

        # Set new rotation
        shape.rotation = new_rotation

        logger.debug(f"Mirrored arrow rotation: {current_rotation}° → {new_rotation}°")

    except Exception as e:
        logger.warning(f"Could not flip arrow rotation: {str(e)}")

def _set_arabic_font(shape) -> None:
    """
    Set Arabic-compatible font for all text runs

    Common Arabic fonts:
    - Arial (universally supported, good rendering)
    - Calibri (modern, clean)
    - Simplified Arabic
    - Traditional Arabic

    Args:
        shape: PowerPoint shape with text_frame
    """
    if not shape.has_text_frame:
        return

    font_name = Config.ARABIC_FONT

    text_frame = shape.text_frame

    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            try:
                # Set font name
                run.font.name = font_name

                # Ensure font applies to complex scripts (Arabic, Hebrew, etc.)
                # Access XML to set both Latin and Complex Script fonts
                rPr = run._element.rPr
                if rPr is not None:
                    # Set Latin font
                    latin_elem = rPr.find('.//a:latin', namespaces=NAMESPACES)
                    if latin_elem is not None:
                        latin_elem.set('typeface', font_name)
                    else:
                        # Create latin element
                        latin_elem = etree.SubElement(rPr, f"{{{NAMESPACES['a']}}}latin")
                        latin_elem.set('typeface', font_name)

                    # Set Complex Script font (for Arabic)
                    cs_elem = rPr.find('.//a:cs', namespaces=NAMESPACES)
                    if cs_elem is not None:
                        cs_elem.set('typeface', font_name)
                    else:
                        # Create cs element
                        cs_elem = etree.SubElement(rPr, f"{{{NAMESPACES['a']}}}cs")
                        cs_elem.set('typeface', font_name)

                logger.debug(f"Set Arabic font: {font_name}")

            except Exception as e:
                logger.warning(f"Could not set font for run: {str(e)}")

def reverse_bullet_order(text_frame) -> None:
    """
    Reverse the order of bullet points (OPTIONAL - not used in MVP)

    Some clients may want bullets visually reversed for RTL.
    For MVP, we skip this as it adds complexity and isn't always required.

    Args:
        text_frame: PowerPoint text_frame object
    """
    # Extract all paragraph data
    paragraphs_data = []

    for paragraph in text_frame.paragraphs:
        paragraphs_data.append({
            'text': paragraph.text,
            'level': paragraph.level,
            'font_name': paragraph.font.name if paragraph.font.name else None,
            'font_size': paragraph.font.size,
            'bold': paragraph.font.bold,
            'italic': paragraph.font.italic
        })

    # Clear existing paragraphs
    text_frame.clear()

    # Re-add in reverse order
    for data in reversed(paragraphs_data):
        p = text_frame.add_paragraph()
        p.text = data['text']
        p.level = data['level']

        if data['font_name']:
            p.font.name = data['font_name']
        if data['font_size']:
            p.font.size = data['font_size']
        if data['bold']:
            p.font.bold = data['bold']
        if data['italic']:
            p.font.italic = data['italic']

    logger.info("Reversed bullet order")

if __name__ == "__main__":
    # Test RTL conversion
    if len(sys.argv) > 2:
        input_file = sys.argv[1]
        output_file = sys.argv[2]
        flip_to_rtl_layout(input_file, output_file)
        print(f"RTL conversion complete: {output_file}")
    else:
        print("Usage: python rtl_converter.py <input.pptx> <output.pptx>")
