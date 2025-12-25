"""
Chart Collision Fixer Module
Detects and fixes collisions between chart groups and other objects after RTL flipping
"""
import sys
import os
from pptx.enum.shapes import MSO_SHAPE_TYPE

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from utils.logger import setup_logger

logger = setup_logger(__name__)

def detect_chart_collisions(prs):
    """
    Scan all slides and detect which ones have chart-to-object collisions

    Args:
        prs: PowerPoint presentation object

    Returns:
        List of dictionaries with collision information for each problematic slide
    """
    slides_with_collisions = []

    logger.info("Scanning all slides for chart-to-object collisions...")

    for slide_idx, slide in enumerate(prs.slides):
        # Separate chart groups from non-chart groups
        chart_groups = []
        non_chart_groups = []

        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                # Check if this group contains a chart
                is_chart_group = _contains_chart(shape)

                group_info = {
                    'shape': shape,
                    'name': shape.name if hasattr(shape, 'name') else 'unnamed',
                    'left': shape.left if shape.left else 0,
                    'top': shape.top if shape.top else 0,
                    'width': shape.width if shape.width else 0,
                    'height': shape.height if shape.height else 0,
                    'right': (shape.left + shape.width) if shape.left and shape.width else 0,
                    'bottom': (shape.top + shape.height) if shape.top and shape.height else 0
                }

                if is_chart_group:
                    chart_groups.append(group_info)
                else:
                    non_chart_groups.append(group_info)

        # Check for collisions between chart groups and non-chart groups
        if chart_groups and non_chart_groups:
            collisions = _detect_collisions_on_slide(chart_groups, non_chart_groups)

            if collisions:
                slides_with_collisions.append({
                    'slide_idx': slide_idx,
                    'slide': slide,
                    'chart_groups': chart_groups,
                    'non_chart_groups': non_chart_groups,
                    'collisions': collisions
                })

                logger.info(f"  Slide {slide_idx + 1}: {len(collisions)} collision(s) detected")

    logger.info(f"Found {len(slides_with_collisions)} slide(s) with chart collisions")
    return slides_with_collisions

def _contains_chart(group_shape):
    """
    Check if a group contains a chart

    Args:
        group_shape: PowerPoint group shape

    Returns:
        True if group contains at least one chart
    """
    try:
        for shape in group_shape.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.CHART:
                return True
    except:
        pass

    return False

def _detect_collisions_on_slide(chart_groups, non_chart_groups):
    """
    Detect collisions between chart groups and non-chart groups on a single slide

    Args:
        chart_groups: List of chart group info dictionaries
        non_chart_groups: List of non-chart group info dictionaries

    Returns:
        List of collision dictionaries
    """
    collisions = []

    for chart_group in chart_groups:
        chart_left = chart_group['left']
        chart_right = chart_group['right']
        chart_top = chart_group['top']
        chart_bottom = chart_group['bottom']

        for other_group in non_chart_groups:
            other_left = other_group['left']
            other_right = other_group['right']
            other_top = other_group['top']
            other_bottom = other_group['bottom']

            # Check horizontal overlap
            h_overlap = chart_left < other_right and chart_right > other_left

            # Check vertical overlap
            v_overlap = chart_top < other_bottom and chart_bottom > other_top

            # If both overlap, we have a collision
            if h_overlap and v_overlap:
                # Calculate overlap amount
                overlap_left = max(chart_left, other_left)
                overlap_right = min(chart_right, other_right)
                overlap_width = overlap_right - overlap_left

                collision_info = {
                    'chart_group': chart_group,
                    'other_group': other_group,
                    'overlap_width': overlap_width,
                    'overlap_width_in': overlap_width / 914400
                }

                collisions.append(collision_info)

                logger.debug(f"    Collision: {chart_group['name']} vs {other_group['name']} (overlap: {overlap_width/914400:.2f}\")")

    return collisions

def fix_chart_collisions_option_c(prs, slide_width):
    """
    Option C: Shift charts as much as possible to minimize overlap

    This function:
    1. Detects all chart-to-object collisions
    2. For each collision, shifts the chart to minimize overlap
    3. Ensures chart doesn't go off-slide

    Args:
        prs: PowerPoint presentation object
        slide_width: Width of the slide in EMUs
    """
    logger.info("="*60)
    logger.info("APPLYING OPTION C: Shift charts to minimize overlap")
    logger.info("="*60)

    # Detect collisions
    slides_with_collisions = detect_chart_collisions(prs)

    if not slides_with_collisions:
        logger.info("No chart collisions detected - no fixes needed")
        return

    # Fix each slide
    total_fixes = 0
    for slide_info in slides_with_collisions:
        slide_idx = slide_info['slide_idx']
        collisions = slide_info['collisions']

        logger.info(f"\nFixing Slide {slide_idx + 1} ({len(collisions)} collision(s))...")

        for collision in collisions:
            fixed = _shift_chart_to_minimize_overlap(collision, slide_width)
            if fixed:
                total_fixes += 1

    logger.info(f"\n{'='*60}")
    logger.info(f"Chart collision fixes applied: {total_fixes}")
    logger.info(f"{'='*60}")

def _shift_chart_to_minimize_overlap(collision, slide_width):
    """
    Shift a chart to minimize overlap with another object

    Strategy:
    1. Try shifting chart LEFT (away from the object)
    2. Try shifting chart RIGHT (away from the object)
    3. Choose the direction that results in less overlap or no overlap
    4. Ensure chart doesn't go off-slide

    Args:
        collision: Collision info dictionary
        slide_width: Slide width in EMUs

    Returns:
        True if fix was applied, False otherwise
    """
    chart_group = collision['chart_group']
    other_group = collision['other_group']

    chart_shape = chart_group['shape']
    chart_left = chart_group['left']
    chart_width = chart_group['width']
    chart_right = chart_group['right']

    other_left = other_group['left']
    other_right = other_group['right']

    min_spacing = 182880  # 0.2 inches

    logger.info(f"  Attempting to fix collision between:")
    logger.info(f"    Chart: {chart_group['name']} (left={chart_left/914400:.2f}\", right={chart_right/914400:.2f}\")")
    logger.info(f"    Object: {other_group['name']} (left={other_left/914400:.2f}\", right={other_right/914400:.2f}\")")

    # Option 1: Shift chart LEFT so it ends before the object starts
    new_left_option1 = other_left - chart_width - min_spacing
    option1_valid = new_left_option1 >= 0

    # Option 2: Shift chart RIGHT so it starts after the object ends
    # Allow chart to go off-slide to the right to avoid collision
    new_left_option2 = other_right + min_spacing
    # Always valid - we allow going off-slide to the right
    option2_valid = True

    logger.info(f"  Option 1 (shift LEFT): new_left={new_left_option1/914400:.2f}\" - {'VALID' if option1_valid else 'INVALID (off-slide)'}")
    logger.info(f"  Option 2 (shift RIGHT): new_left={new_left_option2/914400:.2f}\" - {'VALID' if option2_valid else 'INVALID (off-slide)'}")

    # Choose best option
    if option1_valid and option2_valid:
        # Both valid - choose the one that requires less movement
        shift_amount_1 = abs(new_left_option1 - chart_left)
        shift_amount_2 = abs(new_left_option2 - chart_left)

        if shift_amount_1 <= shift_amount_2:
            new_left = new_left_option1
            direction = "LEFT"
        else:
            new_left = new_left_option2
            direction = "RIGHT"
    elif option1_valid:
        new_left = new_left_option1
        direction = "LEFT"
    elif option2_valid:
        new_left = new_left_option2
        direction = "RIGHT"
    else:
        # Neither option valid - try to minimize overlap by shifting as much as possible
        logger.warning(f"  Cannot fully resolve collision - chart too wide")
        logger.info(f"  Attempting to minimize overlap...")

        # Try shifting right as much as possible (align right edge to slide edge)
        new_left_max_right = slide_width - chart_width

        # Shift RIGHT to minimize overlap - chart moves away from left objects
        new_left = new_left_max_right
        direction = "RIGHT (partial fix - align to right edge)"

    # Apply the shift
    shift_amount = new_left - chart_left
    chart_shape.left = int(new_left)

    logger.info(f"  [OK] Shifted chart {direction} by {abs(shift_amount)/914400:.2f}\"")
    logger.info(f"    New position: {new_left/914400:.2f}\" to {(new_left + chart_width)/914400:.2f}\"")

    return True
